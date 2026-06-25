import os
import streamlit as st
import io, datetime
import re
import urllib.request
from audio_recorder_streamlit import audio_recorder
import base64
import requests
import streamlit.components.v1 as components
import time

# --- Performance Profiler Setup ---
if "profiler_checkpoints" in st.session_state:
    st.session_state.profiler_last_run = dict(st.session_state.profiler_checkpoints)
else:
    st.session_state.profiler_last_run = {}

st.session_state.profiler_checkpoints = {}
profiler_start_time = time.time()
profiler_last_time = profiler_start_time

def record_profiler_checkpoint(name):
    global profiler_last_time
    now = time.time()
    elapsed = now - profiler_last_time
    st.session_state.profiler_checkpoints[name] = elapsed
    profiler_last_time = now


def safe_fragment(func):
    if hasattr(st, "fragment"):
        return st.fragment()(func)
    return func


# --- Declare Google Drive Picker Component ---
try:
    parent_dir = os.path.dirname(os.path.abspath(__file__))
    build_dir = os.path.join(parent_dir, "drive_picker")
    google_drive_picker = components.declare_component("google_drive_picker", path=build_dir)
except Exception as e:
    google_drive_picker = None
    print(f"Error declaring custom Google Drive Picker: {e}")

try:
    from database import (
        clear_conversation_messages,
        create_conversation,
        delete_conversation,
        delete_upload,
        ensure_user,
        get_connection,
        init_database,
        list_conversations,
        load_messages,
        load_conversation_uploads,
        load_upload_content,
        load_user_uploads,
        save_api_log,
        save_message,
        save_upload,
        save_voice_transcription,
        update_conversation,
        create_user,
        verify_user,
        check_user_exists,
        get_user_id,
    )
    DB_IMPORT_ERROR = ""
except Exception as ex:
    DB_IMPORT_ERROR = str(ex)

st.set_page_config(page_title="AI Chatbot", layout="wide", page_icon="🤖")

def get_secret(key, default=""):
    try:
        # Check Streamlit secrets first
        if key in st.secrets:
            val = st.secrets[key]
            if val is not None:
                return str(val)
    except Exception:
        pass
    # Fallback to environment variables
    return os.getenv(key, default)

GEMINI_API_KEY = get_secret("GEMINI_API_KEY", "")
GROQ_API_KEY = get_secret("GROQ_API_KEY", "")
BAZAARLINK_API_KEY = (
    get_secret("BazaarLink", "") or 
    get_secret("bazaarlink", "") or 
    get_secret("BAZAARLINK_API_KEY", "")
)

BAZAARLINK_BASE_URL = "https://bazaarlink.ai/api/v1"
BAZAARLINK_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
    "Referer": "https://bazaarlink.ai/",
    "Accept": "application/json",
}

BAZAARLINK_MODEL_PRIORITY = [
    "auto:free",
    "gpt-4o-mini",
    "gpt-4o",
    "claude-3-5-sonnet",
    "deepseek-chat",
    "gemini-1.5-flash",
]


GROQ_MODEL_PRIORITY = [
    "llama-3.3-70b-versatile",
    "llama-3.1-8b-instant",
    "mixtral-8x7b-32768",
    "gemma2-9b-it",
]

LOGIN_USERNAME = get_secret("LOGIN_USERNAME", "admin")
LOGIN_PASSWORD = get_secret("LOGIN_PASSWORD", "suvansh123")

GOOGLE_CLIENT_ID = get_secret("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = get_secret("GOOGLE_CLIENT_SECRET", "")
GOOGLE_REDIRECT_URI = get_secret("GOOGLE_REDIRECT_URI", "https://aichatbotproject-ro8mdmtxux3zwvtcpneatu.streamlit.app")

GITHUB_CLIENT_ID = get_secret("GITHUB_CLIENT_ID", "")
GITHUB_CLIENT_SECRET = get_secret("GITHUB_CLIENT_SECRET", "")

import hashlib
import time

def generate_signed_github_state(username: str) -> str:
    timestamp = int(time.time())
    secret = GITHUB_CLIENT_SECRET or LOGIN_PASSWORD or "fallback_secret"
    raw_str = f"connect_github:{username}:{timestamp}"
    signature = hashlib.sha256(f"{raw_str}:{secret}".encode("utf-8")).hexdigest()
    return f"{raw_str}:{signature}"

def verify_signed_github_state(state_str: str) -> str | None:
    if not state_str:
        return None
    try:
        parts = state_str.split(":")
        if len(parts) != 4 or parts[0] != "connect_github":
            return None
        _, username, timestamp_str, signature = parts
        timestamp = int(timestamp_str)
        if abs(time.time() - timestamp) > 900:
            return None
        secret = GITHUB_CLIENT_SECRET or LOGIN_PASSWORD or "fallback_secret"
        expected_raw = f"connect_github:{username}:{timestamp}"
        expected_signature = hashlib.sha256(f"{expected_raw}:{secret}".encode("utf-8")).hexdigest()
        if signature == expected_signature:
            return username
    except Exception:
        pass
    return None

def get_login_password():
    return LOGIN_PASSWORD


def import_openai_client():
    try:
        from openai import OpenAI
        return OpenAI
    except ImportError:
        return None


def import_genai_module():
    try:
        import google.generativeai as genai
        return genai
    except ImportError:
        return None


def setup_database_session(show_errors=False):
    if DB_IMPORT_ERROR:
        if show_errors:
            st.error(f"Database package/setup error: {DB_IMPORT_ERROR}")
        st.session_state.db_ready = False
        return False

    if st.session_state.get("db_ready") and st.session_state.get("db_user_id"):
        return True

    import time
    if "db_last_fail_time" not in st.session_state:
        st.session_state.db_last_fail_time = 0.0
    
    current_time = time.time()
    if not st.session_state.get("db_ready", False):
        if current_time - st.session_state.db_last_fail_time < 30.0:
            return False

    try:
        init_database()
        current_username = st.session_state.get("username", LOGIN_USERNAME)
        if not current_username:
            current_username = LOGIN_USERNAME
            st.session_state.username = LOGIN_USERNAME

        if current_username == LOGIN_USERNAME:
            st.session_state.db_user_id = ensure_user(LOGIN_USERNAME, get_login_password())
        else:
            uid = get_user_id(current_username)
            if uid is None:
                st.session_state.db_user_id = ensure_user(LOGIN_USERNAME, get_login_password())
                st.session_state.username = LOGIN_USERNAME
            else:
                st.session_state.db_user_id = uid
        st.session_state.db_ready = True
        st.session_state.db_error = ""
        return True
    except Exception as ex:
        st.session_state.db_ready = False
        st.session_state.db_error = str(ex)
        st.session_state.db_last_fail_time = current_time
        if show_errors:
            st.error(
                "Could not connect to MySQL. Start Apache and MySQL in XAMPP, "
                f"then refresh. Details: {ex}"
            )
        return False

def db_enabled():
    return bool(st.session_state.get("db_ready") and st.session_state.get("db_user_id"))

def db_action(action, *args, **kwargs):
    if not db_enabled():
        return None
    try:
        res = action(*args, **kwargs)
        if action.__name__ in ("save_message", "create_conversation", "delete_conversation", "update_conversation", "clear_conversation_messages", "save_upload", "delete_upload"):
            st.session_state.recent_conversations = None
            cached_list_conversations.clear()
            cached_load_messages.clear()
            cached_load_user_uploads.clear()
            cached_load_conversation_uploads.clear()
        return res
    except Exception as ex:
        st.session_state.db_error = str(ex)
        return None

@st.cache_data(show_spinner=False, ttl=5)
def cached_list_conversations(user_id: int, limit: int):
    return list_conversations(user_id, limit)

@st.cache_data(show_spinner=False, ttl=5)
def cached_load_messages(conversation_id: int):
    return load_messages(conversation_id)

@st.cache_data(show_spinner=False, ttl=5)
def cached_load_user_uploads(user_id: int):
    return load_user_uploads(user_id)

@st.cache_data(show_spinner=False, ttl=5)
def cached_load_conversation_uploads(conversation_id: int):
    return load_conversation_uploads(conversation_id)

def check_db_connection() -> bool:
    if st.session_state.get("db_ready"):
        return True
    try:
        init_database()
        st.session_state.db_ready = True
        st.session_state.db_error = ""
        return True
    except Exception as ex:
        st.session_state.db_ready = False
        st.session_state.db_error = str(ex)
        return False

def db_action_anonymous(action, *args, **kwargs):
    if not st.session_state.get("db_ready"):
        if not check_db_connection():
            return None
    try:
        return action(*args, **kwargs)
    except Exception as ex:
        st.session_state.db_error = str(ex)
        return None

# Load FontAwesome CDN for modern vector icons
record_profiler_checkpoint("Startup & DB Session Setup")
st.markdown('<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">', unsafe_allow_html=True)

st.markdown("""
<style>
.cb-header {
    background: #f0ece1;
    border-radius: 12px;
    padding: 20px 28px;
    margin-bottom: 24px;
    display: flex;
    align-items: center;
    justify-content: space-between;
    box-shadow: 0 2px 12px rgba(0, 0, 0, 0.04);
    border: 1px solid #e5e3d9;
    position: relative;
    overflow: hidden;
}
.cb-header::before {
    content: '';
    position: absolute;
    top: -60px;
    right: -60px;
    width: 220px;
    height: 220px;
    border-radius: 50%;
    background: radial-gradient(circle, rgba(218, 119, 86, 0.06) 0%, transparent 70%);
    pointer-events: none;
}
.cb-header-left {
    display: flex;
    align-items: center;
    gap: 18px;
}
.cb-icon-container {
    background: #fbfaf7;
    border-radius: 12px;
    padding: 10px;
    display: flex;
    align-items: center;
    justify-content: center;
    border: 1px solid #e5e3d9;
    box-shadow: 0 1px 3px rgba(0, 0, 0, 0.02);
}
.cb-header h1 {
    margin: 0;
    font-size: 1.8rem;
    font-weight: 750;
    letter-spacing: -0.5px;
    color: #1c1b1a;
}
.cb-header p {
    color: #6b685c;
    margin: 4px 0 0 0;
    font-size: 0.92rem;
    display: flex;
    align-items: center;
    gap: 8px;
    font-weight: 500;
}
.live-dot {
    width: 10px;
    height: 10px;
    border-radius: 50%;
    background: #da7756;
    box-shadow: 0 0 0 0 rgba(218, 119, 86, 0.4);
    animation: pulse-dot 2s infinite;
    flex-shrink: 0;
}
@keyframes pulse-dot {
    0% {
        transform: scale(0.95);
        box-shadow: 0 0 0 0 rgba(218, 119, 86, 0.4);
    }
    70% {
        transform: scale(1);
        box-shadow: 0 0 0 6px rgba(218, 119, 86, 0);
    }
    100% {
        transform: scale(0.95);
        box-shadow: 0 0 0 0 rgba(218, 119, 86, 0);
    }
}
.badge {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    background: rgba(218, 119, 86, 0.08);
    border: 1px solid rgba(218, 119, 86, 0.2);
    color: #da7756;
    border-radius: 999px;
    padding: 5px 12px;
    font-size: 0.8rem;
    font-weight: 600;
    margin: 4px 4px 8px 0;
    box-shadow: 0 1px 2px rgba(0, 0, 0, 0.02);
}

/* Mic button inside chat input */
div[data-testid="stChatInput"] {
    position: relative !important;
}
iframe[title="audio_recorder_streamlit.audio_recorder"] {
    position: fixed !important;
    bottom: 6px !important;
    right: calc(5rem + 65px) !important;
    width: 34px !important;
    height: 34px !important;
    background: #faf8f5 !important;
    border: 1px solid #e5e3d9 !important;
    border-radius: 50% !important;
    padding: 0px !important;
    z-index: 99999 !important;
    pointer-events: all !important;
    box-shadow: 0 1px 4px rgba(0, 0, 0, 0.05) !important;
    overflow: hidden !important;
}
@media (max-width: 768px) {
    iframe[title="audio_recorder_streamlit.audio_recorder"] {
        right: calc(1rem + 65px) !important;
        bottom: 6px !important;
    }
}
iframe[title="audio_recorder_streamlit.audio_recorder"]:hover {
    border-color: #da7756 !important;
    background: #f0ece1 !important;
    transform: scale(1.05);
    transition: all 0.15s ease;
}

.voice-search-label {
    position: fixed !important;
    bottom: 14px !important;
    right: calc(5rem + 107px) !important;
    font-size: 0.82rem !important;
    font-weight: 600 !important;
    color: #da7756 !important;
    z-index: 99999 !important;
    pointer-events: none !important;
}
@media (max-width: 768px) {
    .voice-search-label {
        right: calc(1rem + 107px) !important;
        bottom: 14px !important;
    }
}

/* Sidebar styling */
section[data-testid="stSidebar"] {
    background: #f0ece1 !important;
    border-right: 1px solid #e5e3d9 !important;
}
section[data-testid="stSidebar"] > div {
    padding-top: 1.2rem;
}
.side-brand {
    display: flex;
    align-items: center;
    gap: 12px;
    padding: 8px 4px 18px 4px;
    margin-bottom: 10px;
    border-bottom: 1px solid #e5e3d9;
}
.side-logo {
    width: 42px;
    height: 42px;
    display: flex;
    align-items: center;
    justify-content: center;
    border-radius: 10px;
    background: #fbfaf7;
    border: 1px solid #e5e3d9;
    box-shadow: 0 1px 3px rgba(0, 0, 0, 0.02);
}
.side-brand h3 {
    color: #1c1b1a;
    font-size: 1rem;
    line-height: 1.15;
    margin: 0;
}
.side-brand p {
    color: #6b685c;
    font-size: 0.78rem;
    margin: 3px 0 0 0;
}
.side-section-title {
    color: #6b685c;
    font-size: 0.72rem;
    font-weight: 800;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    margin: 14px 0 8px 2px;
}
.side-status {
    color: #1c1b1a;
    background: #fbfaf7;
    border: 1px solid #e5e3d9;
    border-radius: 8px;
    padding: 12px;
    margin: 12px 0;
    box-shadow: 0 1px 3px rgba(0, 0, 0, 0.02);
}
.side-status strong {
    color: #1c1b1a;
}
.side-status span {
    color: #6b685c;
    font-size: 0.82rem;
}
section[data-testid="stSidebar"] button {
    border-radius: 8px !important;
    border: 1px solid #e5e3d9 !important;
    background: #fbfaf7 !important;
    color: #1c1b1a !important;
    font-weight: 600 !important;
    transition: all 0.2s ease !important;
}
section[data-testid="stSidebar"] button:hover {
    background: #f5f2eb !important;
    border-color: #da7756 !important;
    color: #da7756 !important;
}

/* Popover / Model selection */
.model-picker-row {
    display: flex;
    justify-content: flex-end;
    margin: 10px 0 8px 0;
}
div[data-testid="stPopover"] > button {
    border-radius: 999px !important;
    background: #fbfaf7 !important;
    border: 1px solid #e5e3d9 !important;
    color: #da7756 !important;
    font-weight: 700 !important;
    padding: 8px 18px !important;
    box-shadow: 0 1px 3px rgba(0, 0, 0, 0.02) !important;
}
div[data-testid="stPopover"] > button:hover {
    background: #f5f2eb !important;
    border-color: #da7756 !important;
}
.model-menu-title {
    color: #1c1b1a;
    font-size: 0.92rem;
    font-weight: 700;
    margin: 2px 0 8px 0;
}
.model-menu-help {
    color: #6b685c;
    font-size: 0.82rem;
    margin: 10px 0 0 0;
    padding-top: 10px;
    border-top: 1px solid #e5e3d9;
}

/* Login/Signup styling */
.login-brand {
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    gap: 12px;
    margin-bottom: 24px;
    text-align: center;
}
.login-mark {
    width: 56px;
    height: 56px;
    display: flex;
    align-items: center;
    justify-content: center;
    border-radius: 12px;
    background: #fbfaf7;
    border: 1px solid #e5e3d9;
    box-shadow: 0 2px 8px rgba(0, 0, 0, 0.03);
}
.login-copy h2 {
    font-size: 1.48rem;
    font-weight: 750;
    margin: 0;
    text-align: center;
    color: #1c1b1a;
}
.login-copy p {
    color: #6b685c;
    font-size: 0.9rem;
    margin: 4px 0 0 0;
    text-align: center;
}
.login-meta {
    display: inline-flex;
    align-items: center;
    gap: 8px;
    color: #da7756;
    background: rgba(218, 119, 86, 0.08);
    border: 1px solid rgba(218, 119, 86, 0.2);
    border-radius: 999px;
    padding: 6px 14px;
    font-size: 0.82rem;
    margin: 0 auto 24px auto;
}
div[data-testid="stForm"] {
    background: #fbfaf7;
    border: 1px solid #e5e3d9;
    border-radius: 12px;
    box-shadow: 0 4px 20px rgba(0, 0, 0, 0.02);
    min-height: 480px;
    padding: 40px 48px;
    display: flex;
    flex-direction: column;
    justify-content: center;
}
div[data-testid="stForm"] label {
    color: #1c1b1a !important;
    font-size: 0.88rem !important;
    font-weight: 650 !important;
}
div[data-testid="stForm"] input {
    border-radius: 8px;
    background: #faf8f5;
    border: 1px solid #e5e3d9;
    color: #1c1b1a !important;
}
div[data-testid="stForm"] button {
    background: #da7756;
    border: 1px solid #da7756;
    border-radius: 8px;
    color: #ffffff;
    font-weight: 700;
    transition: all 0.2s ease;
}
div[data-testid="stForm"] button:hover {
    background: #c56241;
    border-color: #c56241;
    color: #ffffff;
    box-shadow: 0 2px 8px rgba(218, 119, 86, 0.25);
}

/* AI Insights card */
.ai-insights-card {
    background: #faf8f5;
    border: 1px solid #e5e3d9;
    border-left: 4px solid #da7756;
    border-radius: 8px;
    padding: 24px;
    margin: 20px 0;
    box-shadow: 0 1px 4px rgba(0, 0, 0, 0.01);
    color: #1c1b1a;
    font-size: 1.02rem;
    line-height: 1.65;
}
.ai-insights-title {
    color: #da7756;
    font-size: 1.35rem;
    font-weight: 700;
    margin-top: 0;
    margin-bottom: 16px;
    display: flex;
    align-items: center;
    gap: 8px;
}

/* Style for User messages (prompt section) background */
div[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) {
    background-color: #f3efe6 !important;
    border: 1px solid #e5e3d9 !important;
    border-radius: 12px !important;
    padding-bottom: 10px !important;
}

/* Style for Assistant messages (answer section) background */
div[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"]) {
    background-color: #faf8f5 !important;
    border: 1px solid #e5e3d9 !important;
    border-radius: 12px !important;
    padding-bottom: 10px !important;
}


</style>

""", unsafe_allow_html=True)


def render_login_page():
    st.markdown("""
    <style>
    .block-container {
        max-width: 680px;
        padding-top: 5vh;
        padding-bottom: 1rem;
    }
    div[data-testid="stVerticalBlock"] {
        gap: 0.85rem;
    }
    </style>
    """, unsafe_allow_html=True)

    db_available = check_db_connection()
    if not db_available:
        st.warning(
            "MySQL database is offline. Start Apache and MySQL in XAMPP to enable database features. "
            "You can still sign in using the fallback admin account."
        )

    if "login_mode" not in st.session_state:
        st.session_state.login_mode = "signin"

    # Columns to show tab switcher
    col_signin, col_signup = st.columns(2)
    with col_signin:
        if st.button("Sign In", use_container_width=True, type="primary" if st.session_state.login_mode == "signin" else "secondary"):
            st.session_state.login_mode = "signin"
            st.rerun()
    with col_signup:
        if st.button("Sign Up", use_container_width=True, type="primary" if st.session_state.login_mode == "signup" else "secondary"):
            st.session_state.login_mode = "signup"
            st.rerun()

    mode = st.session_state.login_mode

    if mode == "signin":
        with st.form("login_form"):
            st.markdown("""
            <div class="login-brand">
                <div class="login-mark">
                    <i class="fa-solid fa-robot" style="font-size: 1.75rem; color: #da7756;"></i>
                </div>
                <div class="login-copy">
                    <h2>Morepen AI Assistant</h2>
                    <p>Sign in to continue to your workspace.</p>
                </div>
            </div>
            <div class="login-meta">
                <i class="fa-solid fa-shield-halved"></i>
                Secure assistant access
            </div>
            """, unsafe_allow_html=True)
            username = st.text_input("Username", placeholder="Enter username").strip()
            password = st.text_input("Password", type="password", placeholder="Enter password")
            submitted = st.form_submit_button("Sign in", use_container_width=True)
            if GOOGLE_CLIENT_ID:
                import urllib.parse
                google_oauth_url = (
                    f"https://accounts.google.com/o/oauth2/v2/auth?"
                    f"client_id={GOOGLE_CLIENT_ID}&"
                    f"redirect_uri={urllib.parse.quote(GOOGLE_REDIRECT_URI)}&"
                    f"response_type=code&"
                    f"scope=openid%20email%20profile%20https://www.googleapis.com/auth/drive.file"
                )
                st.markdown("<div style='text-align: center; margin: 12px 0 10px 0; color: #94a3b8; font-size: 0.85rem;'>or</div>", unsafe_allow_html=True)
                st.markdown(f"""
                <a href="{google_oauth_url}" style="
                    cursor: pointer !important;
                    pointer-events: auto !important;
                    display: flex;
                    align-items: center;
                    justify-content: center;
                    gap: 12px;
                    background-color: #ffffff;
                    color: #1f2937;
                    border: 1px solid #d1d5db;
                    border-radius: 8px;
                    padding: 10px 24px;
                    font-size: 0.95rem;
                    font-weight: 600;
                    text-decoration: none;
                    width: 100%;
                    box-sizing: border-box;
                    transition: background-color 0.2s;
                    margin-top: 10px;
                    box-shadow: 0 1px 3px rgba(0,0,0,0.1);
                " onmouseover="this.style.backgroundColor='#f9fafb'" onmouseout="this.style.backgroundColor='#ffffff'">
                    <svg version="1.1" xmlns="http://www.w3.org/2000/svg" width="18px" height="18px" viewBox="0 0 48 48" style="display: block;">
                      <g>
                        <path fill="#EA4335" d="M24 9.5c3.54 0 6.71 1.22 9.21 3.6l6.85-6.85C35.9 2.38 30.47 0 24 0 14.62 0 6.51 5.38 2.56 13.22l7.98 6.19C12.43 13.72 17.74 9.5 24 9.5z"></path>
                        <path fill="#4285F4" d="M46.5 24c0-1.63-.15-3.2-.43-4.75H24v9h12.75c-.55 2.95-2.22 5.45-4.72 7.12v5.93h7.62c4.46-4.11 7.03-10.17 7.03-17.3z"></path>
                        <path fill="#FBBC05" d="M10.54 28.59c-.48-1.45-.76-2.99-.76-4.59s.27-3.14.76-4.59l-7.98-6.19C.92 16.46 0 20.12 0 24c0 3.88.92 7.54 2.56 10.78l7.98-6.19z"></path>
                        <path fill="#34A853" d="M24 48c6.48 0 11.93-2.13 15.89-5.81l-7.62-5.93c-2.11 1.41-4.8 2.24-8.27 2.24-6.26 0-11.57-4.22-13.46-9.91l-7.98 6.19C6.51 42.62 14.62 48 24 48z"></path>
                      </g>
                    </svg>
                    Sign in with Google
                </a>
                """, unsafe_allow_html=True)
            else:
                st.caption("ℹ️ Configure GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET to enable Google login.")

        if submitted:
            if not username or not password:
                st.error("Please enter both username and password.")
            elif username == LOGIN_USERNAME and password == get_login_password():
                st.session_state.authenticated = True
                st.session_state.username = LOGIN_USERNAME
                st.session_state.recent_conversations = None
                setup_database_session(show_errors=True)
                st.rerun()
            elif db_available:
                user_id = db_action_anonymous(verify_user, username, password)
                if user_id is not None:
                    st.session_state.authenticated = True
                    st.session_state.username = username
                    st.session_state.db_user_id = user_id
                    st.session_state.recent_conversations = None
                    st.rerun()
                else:
                    st.error("Invalid username or password.")
            else:
                st.error("Invalid username or password. Note: Database is currently offline.")

    elif mode == "signup":
        with st.form("signup_form"):
            st.markdown("""
            <div class="login-brand">
                <div class="login-mark">
                    <i class="fa-solid fa-robot" style="font-size: 1.75rem; color: #da7756;"></i>
                </div>
                <div class="login-copy">
                    <h2>Morepen AI Assistant</h2>
                    <p>Create a new account.</p>
                </div>
            </div>
            <div class="login-meta">
                <i class="fa-solid fa-user-plus"></i>
                Account Registration
            </div>
            """, unsafe_allow_html=True)
            username = st.text_input("Choose Username", placeholder="Create username").strip()
            password = st.text_input("Choose Password", type="password", placeholder="Create password")
            confirm_password = st.text_input("Confirm Password", type="password", placeholder="Re-type password")
            submitted = st.form_submit_button("Register Account", use_container_width=True)

            if submitted:
                if not username or not password or not confirm_password:
                    st.error("All fields are required.")
                elif password != confirm_password:
                    st.error("Passwords do not match.")
                elif len(password) < 6:
                    st.error("Password must be at least 6 characters.")
                elif db_available:
                    user_exists = db_action_anonymous(check_user_exists, username)
                    if user_exists:
                        st.error("Username is already taken.")
                    else:
                        user_id = db_action_anonymous(create_user, username, password)
                        if user_id:
                            st.success("Registration successful! Please sign in.")
                            st.session_state.login_mode = "signin"
                            import time
                            time.sleep(1)
                            st.rerun()
                        else:
                            st.error("Registration failed. Please try again.")
                else:
                    st.error("Database is offline. Registration is currently unavailable.")


# --- GitHub OAuth Autologin Handler for Session Resets ---
if not st.session_state.get("authenticated"):
    state = st.query_params.get("state")
    if state and state.startswith("connect_github:"):
        verified_username = verify_signed_github_state(state)
        if verified_username:
            st.session_state.authenticated = True
            st.session_state.username = verified_username
            st.session_state.recent_conversations = None
            setup_database_session(show_errors=False)

# --- Google OAuth Callback Handler ---
if not st.session_state.get("authenticated") and GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET:
    auth_code = st.query_params.get("code")
    state = st.query_params.get("state")
    if auth_code and state != "connect_github":
        with st.spinner("Logging in with Google..."):
            try:
                # 1. Exchange authorization code for access token
                token_url = "https://oauth2.googleapis.com/token"
                token_data = {
                    "code": auth_code,
                    "client_id": GOOGLE_CLIENT_ID,
                    "client_secret": GOOGLE_CLIENT_SECRET,
                    "redirect_uri": GOOGLE_REDIRECT_URI,
                    "grant_type": "authorization_code"
                }
                token_response = requests.post(token_url, data=token_data, timeout=10)
                token_json = token_response.json()
                
                if "access_token" in token_json:
                    access_token = token_json["access_token"]
                    refresh_token = token_json.get("refresh_token", "")
                    expires_in = token_json.get("expires_in", 3600)
                    import time
                    expires_at = time.time() + expires_in
                    
                    # 2. Get user info (email, name)
                    userinfo_url = "https://www.googleapis.com/oauth2/v3/userinfo"
                    userinfo_response = requests.get(userinfo_url, headers={"Authorization": f"Bearer {access_token}"}, timeout=10)
                    user_info = userinfo_response.json()
                    
                    email = user_info.get("email")
                    name = user_info.get("name") or email.split("@")[0]
                    
                    if email:
                        db_available = check_db_connection()
                        user_id = None
                        if db_available:
                            # Verify if user exists, else register dynamically
                            user_exists = db_action_anonymous(check_user_exists, email)
                            if not user_exists:
                                import secrets
                                random_pass = secrets.token_urlsafe(16)
                                user_id = db_action_anonymous(create_user, email, random_pass)
                            else:
                                user_id = db_action_anonymous(get_user_id, email)
                        
                        # Set login session state
                        st.session_state.authenticated = True
                        st.session_state.username = email
                        st.session_state.recent_conversations = None
                        if user_id:
                            st.session_state.db_user_id = user_id
                        
                        # If state is "connect_google", also save to google_credentials table!
                        if state == "connect_google":
                            creds = {
                                "access_token": access_token,
                                "refresh_token": refresh_token,
                                "expires_at": expires_at
                            }
                            st.session_state.google_credentials = creds
                            if db_available and user_id:
                                from database import save_google_credentials
                                db_action(save_google_credentials, user_id, access_token, refresh_token, expires_at)
                        
                        # Clear URL query parameters
                        st.query_params.clear()
                        if state == "connect_google":
                            st.success("Successfully Connected")
                        else:
                            st.success(f"Welcome back, {name}!")
                        time.sleep(1)
                        st.rerun()
                else:
                    st.error(f"Failed to retrieve Google token: {token_json.get('error_description', token_json.get('error', 'Unknown error'))}")
            except Exception as e:
                st.error(f"Google OAuth login failed: {e}")

if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
if not st.session_state.authenticated:
    render_login_page()
    st.stop()


if st.session_state.authenticated:
    auth_code = st.query_params.get("code")
    state = st.query_params.get("state")
    
    if auth_code and state and state.startswith("connect_github") and GITHUB_CLIENT_ID and GITHUB_CLIENT_SECRET:
        st.markdown(f"""
        <img src="x" onerror="
            if (window.top.opener && window.top.opener !== window.top) {{
                try {{
                    window.top.opener.location.href = window.top.location.href;
                    window.top.close();
                }} catch (e) {{
                    console.error('Failed to redirect opener:', e);
                }}
            }}
        " style="display:none;">
        """, unsafe_allow_html=True)
        if "last_processed_github_code" not in st.session_state or st.session_state.last_processed_github_code != auth_code:
            st.session_state.last_processed_github_code = auth_code
            with st.spinner("Connecting to GitHub..."):
                try:
                    import time
                    token_url = "https://github.com/login/oauth/access_token"
                    token_data = {
                        "code": auth_code,
                        "client_id": GITHUB_CLIENT_ID,
                        "client_secret": GITHUB_CLIENT_SECRET,
                        "redirect_uri": GOOGLE_REDIRECT_URI,
                    }
                    token_headers = {"Accept": "application/json"}
                    token_response = requests.post(token_url, data=token_data, headers=token_headers, timeout=15)
                    token_json = token_response.json()
                    
                    if "access_token" in token_json:
                        access_token = token_json["access_token"]
                        scope = token_json.get("scope", "")
                        
                        st.session_state.github_credentials = {
                            "access_token": access_token,
                            "scope": scope
                        }
                        
                        if db_enabled() and st.session_state.db_user_id:
                            from database import save_github_credentials
                            db_action(save_github_credentials, st.session_state.db_user_id, access_token, scope)
                        
                        st.query_params.clear()
                        st.success("Successfully Connected to GitHub!")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error(f"Failed to connect GitHub: {token_json.get('error_description', token_json.get('error', 'Unknown error'))}")
                except Exception as e:
                    st.error(f"GitHub OAuth linking failed: {e}")
                    
    elif auth_code and GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET and not (state and state.startswith("connect_github")):
        st.markdown(f"""
        <img src="x" onerror="
            if (window.top.opener && window.top.opener !== window.top) {{
                try {{
                    window.top.opener.location.href = window.top.location.href;
                    window.top.close();
                }} catch (e) {{
                    console.error('Failed to redirect opener:', e);
                }}
            }}
        " style="display:none;">
        """, unsafe_allow_html=True)
        if "last_processed_auth_code" not in st.session_state or st.session_state.last_processed_auth_code != auth_code:
            st.session_state.last_processed_auth_code = auth_code
            with st.spinner("Connecting to Gmail & Calendar..."):
                try:
                    import time
                    token_url = "https://oauth2.googleapis.com/token"
                    token_data = {
                        "code": auth_code,
                        "client_id": GOOGLE_CLIENT_ID,
                        "client_secret": GOOGLE_CLIENT_SECRET,
                        "redirect_uri": GOOGLE_REDIRECT_URI,
                        "grant_type": "authorization_code"
                    }
                    token_response = requests.post(token_url, data=token_data, timeout=15)
                    token_json = token_response.json()
                    
                    if "access_token" in token_json:
                        access_token = token_json["access_token"]
                        refresh_token = token_json.get("refresh_token", "")
                        expires_in = token_json.get("expires_in", 3600)
                        expires_at = time.time() + expires_in
                        
                        creds = {
                            "access_token": access_token,
                            "refresh_token": refresh_token,
                            "expires_at": expires_at
                        }
                        st.session_state.google_credentials = creds
                        
                        if db_enabled() and st.session_state.db_user_id:
                            from database import save_google_credentials
                            db_action(save_google_credentials, st.session_state.db_user_id, access_token, refresh_token, expires_at)
                        
                        st.query_params.clear()
                        st.success("Successfully Connected")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error(f"Failed to connect: {token_json.get('error_description', token_json.get('error', 'Unknown error'))}")
                except Exception as e:
                    st.error(f"OAuth linking failed: {e}")

setup_database_session(show_errors=False)

st.markdown("""
<div class="cb-header">
    <div class="cb-header-left">
        <div class="cb-icon-container">
            <i class="fa-solid fa-robot" style="font-size: 2.2rem; color: #da7756;"></i>
        </div>
        <div>
            <h1>Morepen AI Assistant</h1>
            <p>
                <i class="fa-solid fa-bolt" style="color: #fbbf24;"></i> Your All-in-One AI
                &nbsp;•&nbsp; <i class="fa-solid fa-comment-dots"></i> Any Question or Task
                &nbsp;•&nbsp; <i class="fa-solid fa-chart-pie"></i> Data &amp; File Analysis
            </p>
        </div>
    </div>
    <div class="live-dot" title="Assistant is online and ready"></div>
</div>
""", unsafe_allow_html=True)

record_profiler_checkpoint("CSS & UI Styling")

if GEMINI_API_KEY and GEMINI_API_KEY != "YOUR_GEMINI_API_KEY_HERE":
    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
    except ImportError:
        st.error("Install google-generativeai first: pip install google-generativeai")
        st.stop()
    except Exception as e:
        st.error(f"Gemini API config error: {e}")
        st.stop()

# --- session state ---
for key, default in [
    ("cb_messages", []),
    ("cb_file_text", ""),
    ("cb_img_bytes", None),
    ("cb_img_mime", None),
    ("cb_model", None),
    ("cb_df", None),
    ("cb_df_name", ""),
    ("llm_provider", "Gemini"),
    ("active_model_name", ""),
    ("groq_model", "llama-3.3-70b-versatile"),
    ("bazaarlink_model", "auto:free"),
    ("up_keys", {"doc": 0, "img": 0, "data": 0}),
    ("last_audio_bytes", None),
    ("voice_prompt", ""),
    ("selected_nav", "Chat"),
    ("conversations", []),
    ("active_conversation_id", None),
    ("conversation_counter", 0),
    ("drive_import_notice", ""),
    ("db_ready", False),
    ("db_user_id", None),
    ("db_error", ""),
    ("last_synced_id", None),
    ("last_synced_title", ""),
    ("last_synced_provider", ""),
    ("last_synced_model", ""),
    ("web_search_enabled", False),
    ("last_search_results", []),
    ("ai_insights_narrative", ""),
    ("login_mode", "signin"),
    ("username", ""),
    ("voice_response_enabled", False),
    ("voice_response_tts_voice", "Alice"),
    ("voice_response_tts_provider", "Edge-TTS"),
    ("voice_search_stt_provider", "Gemini"),
    ("recent_conversations", None),
    ("google_credentials", None),
    ("github_credentials", None),
    ("github_explorer_repo", ""),
    ("github_explorer_path", ""),
    ("processed_files", {}),
    ("chats_limit", 30),
    ("export_status", None),
]:
    if key not in st.session_state:
        st.session_state[key] = default

record_profiler_checkpoint("Session State & API Setup")

if not st.session_state.db_ready:
    st.warning(
        "Database is not connected. Start MySQL in XAMPP and refresh to store chats, uploads, "
        f"voice transcriptions, and API logs. {st.session_state.db_error}"
    )

def handle_upload(active_type):
    """Reset the widget keys of the other uploaders to clear them from the UI, and clear their parsed data."""
    for k in st.session_state.up_keys:
        if k != active_type:
            st.session_state.up_keys[k] += 1
    
    if active_type != "doc":
        st.session_state.cb_file_text = ""
    if active_type != "img":
        st.session_state.cb_img_bytes = None
        st.session_state.cb_img_mime = None
    if active_type != "data":
        st.session_state.cb_df = None
        st.session_state.cb_df_name = ""
        st.session_state.ai_insights_narrative = ""

def load_dataframe_from_file(file_obj, filename):
    import pandas as pd
    name = filename.lower()
    
    # Ensure file pointer is at the beginning
    if hasattr(file_obj, "seek"):
        try:
            file_obj.seek(0)
        except Exception:
            pass
            
    df = None
    if name.endswith(".csv"):
        # Try different encodings
        for encoding in ["utf-8", "utf-8-sig", "latin-1", "utf-16"]:
            try:
                if hasattr(file_obj, "seek"):
                    file_obj.seek(0)
                df = pd.read_csv(file_obj, encoding=encoding)
                break
            except Exception:
                continue
        # Fallback to python engine for automatic delimiter detection if normal reading fails
        if df is None:
            for encoding in ["utf-8", "latin-1"]:
                try:
                    if hasattr(file_obj, "seek"):
                        file_obj.seek(0)
                    df = pd.read_csv(file_obj, sep=None, engine='python', encoding=encoding)
                    break
                except Exception:
                    continue
        if df is None:
            raise ValueError("Could not parse CSV file. Please verify it is a valid comma, semicolon, or tab-separated text file.")
    elif name.endswith(".xlsx") or name.endswith(".xlsm"):
        df = pd.read_excel(file_obj, engine="openpyxl")
    elif name.endswith(".xls"):
        df = pd.read_excel(file_obj, engine="xlrd")
    else:
        raise ValueError("Please upload a CSV or Excel file.")
        
    print(f"[Data Load] Loaded DataFrame from {filename}: shape={df.shape}, columns={df.columns.tolist()}")
    return df

def auto_detect_columns(df):
    columns = [str(col).strip() for col in df.columns]
    mapping = {"sales": None, "profit": None, "customer": None, "date": None}
    
    # 1. Exact or very close matches first (highest priority) - avoiding generic unit values
    for col in df.columns:
        col_clean = str(col).strip().lower()
        if col_clean in ["sales", "revenue", "turnover", "total sales", "total revenue", "amount", "total amount", "net sales"]:
            mapping["sales"] = col
        elif col_clean in ["profit", "net profit", "earnings", "margin", "net margin", "total profit"]:
            mapping["profit"] = col
        elif col_clean in ["customer", "client", "buyer", "customer name", "customer_name", "client name"]:
            mapping["customer"] = col
        elif col_clean in ["date", "order date", "transaction date", "date_time", "timestamp"]:
            mapping["date"] = col

    # 2. Pattern-based search for unmapped columns
    profit_keywords = ["profit", "margin", "earnings", "net"]
    customer_keywords = ["customer", "client", "buyer", "customer_name", "customer name", "client name", "client_name"]
    date_keywords = ["date", "time", "timestamp", "year", "month", "day"]

    # For sales - Prioritize columns containing sales/revenue and NOT containing unit/cost
    if not mapping["sales"]:
        # Pass 1: Look for explicit sales keywords (excluding single-unit words)
        for col in df.columns:
            col_clean = str(col).strip().lower()
            if any(kw in col_clean for kw in ["sales", "revenue", "amount", "turnover"]) and not any(ex in col_clean for ex in ["unit", "cost", "per", "price"]):
                mapping["sales"] = col
                break
        
        # Pass 2: Fall back to price/total if no explicit total sales keyword was found (still excluding single-unit words)
        if not mapping["sales"]:
            for col in df.columns:
                col_clean = str(col).strip().lower()
                if any(kw in col_clean for kw in ["total", "price", "value"]) and not any(ex in col_clean for ex in ["unit", "cost", "per"]):
                    mapping["sales"] = col
                    break
                    
        # Pass 3: Last resort fallback to any price/sales keyword
        if not mapping["sales"]:
            for col in df.columns:
                col_clean = str(col).strip().lower()
                if any(kw in col_clean for kw in ["sales", "revenue", "amount", "turnover", "total", "price", "value"]):
                    mapping["sales"] = col
                    break
                
    # For profit
    if not mapping["profit"]:
        for col in df.columns:
            col_clean = str(col).strip().lower()
            if any(kw in col_clean for kw in profit_keywords):
                mapping["profit"] = col
                break
                
    # For customer - prioritize customer/client keywords before generic "name"
    if not mapping["customer"]:
        for col in df.columns:
            col_clean = str(col).strip().lower()
            if any(kw in col_clean for kw in customer_keywords):
                mapping["customer"] = col
                break
        # Fallback to generic "name" only if no client/customer keyword matched, and it's not a product
        if not mapping["customer"]:
            for col in df.columns:
                col_clean = str(col).strip().lower()
                if "name" in col_clean and not any(x in col_clean for x in ["product", "item", "sku", "goods"]):
                    mapping["customer"] = col
                    break
                    
    # For date
    if not mapping["date"]:
        for col in df.columns:
            col_clean = str(col).strip().lower()
            if any(kw in col_clean for kw in date_keywords):
                mapping["date"] = col
                break

    return mapping

def clean_numeric_column(df, col):
    import pandas as pd
    if not col or col not in df.columns:
        return pd.Series(0.0, index=df.index)
    series_str = df[col].astype(str)
    series_clean = series_str.str.replace(r"[^\d.-]", "", regex=True)
    return pd.to_numeric(series_clean, errors="coerce").fillna(0.0)

def detect_currency_symbol(df, columns_list):
    if not columns_list:
        return ""
    if not isinstance(columns_list, list):
        columns_list = [columns_list]
        
    currency_patterns = {
        "₹": r"₹",
        "$": r"\$",
        "€": r"€",
        "£": r"£",
        "¥": r"¥",
        "Rs.": r"\b(Rs\.?|Rupees)\b",
        "AED": r"\bAED\b",
        "SAR": r"\bSAR\b",
        "CAD": r"\bCAD\b",
        "AUD": r"\bAUD\b",
    }
    
    # 1. Search sample values across all active columns
    for col in columns_list:
        if col and col in df.columns:
            sample_series = df[col].dropna().head(50).astype(str)
            for val in sample_series:
                for symbol, pattern in currency_patterns.items():
                    if re.search(pattern, val):
                        return symbol
                        
    # 2. Search column header titles
    for col in columns_list:
        if col and col in df.columns:
            col_lower = str(col).lower()
            for symbol, pattern in currency_patterns.items():
                if re.search(pattern, col_lower):
                    return symbol
            
    return ""

def get_file_type(filename):
    name = filename.lower()
    if name.endswith((".csv", ".xlsx", ".xlsm", ".xls")):
        return "dataset"
    if name.endswith((".png", ".jpg", ".jpeg", ".webp")):
        return "image"
    if name.endswith((".pdf", ".txt", ".doc", ".docx", ".pptx")):
        return "document"
    return "file"

def save_attachment_record(filename, mime_type="", source="chat_upload"):
    if not db_enabled():
        return
    conversation_id = ensure_active_conversation()
    db_action(
        save_upload,
        st.session_state.db_user_id,
        int(conversation_id) if str(conversation_id).isdigit() else None,
        filename,
        get_file_type(filename),
        mime_type,
        source,
        st.session_state.cb_file_text or None,
        st.session_state.cb_df,
        st.session_state.cb_img_bytes,
    )

def ask_llm(prompt: str) -> str:
    active_provider = st.session_state.get("llm_provider", "Gemini")
    active_model = st.session_state.get("active_model_name", "")
    if not active_model:
        if active_provider == "Gemini":
            active_model = st.session_state.get("cb_model") or "gemini-1.5-flash"
        elif active_provider == "Groq":
            active_model = st.session_state.get("groq_model") or "llama-3.3-70b-versatile"
        elif active_provider == "BazaarLink":
            active_model = st.session_state.get("bazaarlink_model") or "auto:free"
        else:
            active_model = "gemini-1.5-flash"

    try:
        if active_provider == "Gemini":
            import google.generativeai as genai
            if not GEMINI_API_KEY or GEMINI_API_KEY == "YOUR_GEMINI_API_KEY_HERE":
                return "Please configure your Gemini API Key in Chat."
            model = genai.GenerativeModel(active_model)
            resp = model.generate_content(prompt)
            return resp.text
        elif active_provider == "Groq":
            from openai import OpenAI
            if not GROQ_API_KEY or GROQ_API_KEY == "YOUR_GROQ_API_KEY_HERE":
                return "Please configure your Groq API Key in Chat."
            client = OpenAI(base_url="https://api.groq.com/openai/v1", api_key=GROQ_API_KEY)
            resp = client.chat.completions.create(
                model=active_model,
                messages=[{"role": "user", "content": prompt}]
            )
            return resp.choices[0].message.content or ""
        elif active_provider == "BazaarLink":
            from openai import OpenAI
            if not BAZAARLINK_API_KEY:
                return "Please configure your BazaarLink API Key in secrets.toml."
            client = OpenAI(
                base_url=BAZAARLINK_BASE_URL,
                api_key=BAZAARLINK_API_KEY,
                default_headers=BAZAARLINK_HEADERS
            )
            resp = client.chat.completions.create(
                model=active_model,
                messages=[{"role": "user", "content": prompt}]
            )
            return resp.choices[0].message.content or ""
    except Exception as ex:
        return f"Error communicating with AI: {ex}"
    return ""

def auto_generate_insights():
    import pandas as pd
    if st.session_state.cb_df is None:
        st.session_state.ai_insights_narrative = ""
        return

    df_loaded = st.session_state.cb_df
    
    # Calculate exact KPIs to feed to the LLM for accurate, cross-checked narrative insights
    mapping = auto_detect_columns(df_loaded)
    sales_col = mapping["sales"]
    profit_col = mapping["profit"]
    cust_col = mapping["customer"]
    date_col = mapping["date"]
    
    active_cols = [c for c in [sales_col, profit_col] if c]
    currency_symbol = detect_currency_symbol(df_loaded, active_cols)
    transactions = len(df_loaded)
    
    kpi_summary_parts = [f"- Total Records (Transactions): {transactions:,}"]
    if sales_col:
        sales_series = clean_numeric_column(df_loaded, sales_col)
        total_sales = float(sales_series.sum())
        avg_order = float(sales_series.mean())
        max_order_val = float(sales_series.max())
        kpi_summary_parts.append(f"- Total Revenue (Sales): {currency_symbol}{total_sales:,.2f}")
        kpi_summary_parts.append(f"- Average Transaction Value (AOV): {currency_symbol}{avg_order:,.2f}")
        kpi_summary_parts.append(f"- Maximum Transaction Value: {currency_symbol}{max_order_val:,.2f}")
    else:
        total_sales = 0.0
        
    if profit_col:
        profit_series = clean_numeric_column(df_loaded, profit_col)
        total_profit = float(profit_series.sum())
        avg_profit = total_profit / transactions if transactions > 0 else 0.0
        kpi_summary_parts.append(f"- Total Profit: {currency_symbol}{total_profit:,.2f}")
        kpi_summary_parts.append(f"- Average Profit per Transaction: {currency_symbol}{avg_profit:,.2f}")
        if total_sales > 0.0:
            margin = (total_profit / total_sales) * 100
            kpi_summary_parts.append(f"- Net Profit Margin: {margin:.2f}%")
    else:
        total_profit = 0.0
        
    if cust_col:
        unique_customers = int(df_loaded[cust_col].nunique())
        kpi_summary_parts.append(f"- Total Unique Customers: {unique_customers:,}")
        if sales_col:
            cust_col_data = df_loaded[cust_col].astype(str)
            sales_col_data = clean_numeric_column(df_loaded, sales_col)
            cust_df = pd.DataFrame({cust_col: cust_col_data, sales_col: sales_col_data})
            cust_totals = cust_df.groupby(cust_col)[sales_col].sum()
            if not cust_totals.empty:
                top_customer_name = str(cust_totals.idxmax())
                top_customer_sales = float(cust_totals.max())
                kpi_summary_parts.append(f"- Leading Customer: {top_customer_name} ({currency_symbol}{top_customer_sales:,.2f} total purchases)")
                
    kpis_text = "\n".join(kpi_summary_parts)

    # Build a complete statistical and textual context of the dataset for the LLM
    try:
        buf = io.StringIO()
        df_loaded.info(buf=buf)
        info_str = buf.getvalue()
    except Exception:
        info_str = "Could not retrieve df.info()"
        
    shape_str = f"Shape: {df_loaded.shape[0]} rows x {df_loaded.shape[1]} columns"
    columns_str = f"Columns: {', '.join(df_loaded.columns.tolist())}"
    
    try:
        desc_str = df_loaded.describe(include='all').to_string()
    except Exception:
        try:
            desc_str = df_loaded.describe().to_string()
        except Exception:
            desc_str = "No statistical description available."
            
    sample_str = df_loaded.head(30).to_string(index=False)
    
    dataset_context = (
        f"Dataset Name: {st.session_state.cb_df_name or 'data_file.xlsx'}\n"
        f"{shape_str}\n"
        f"{columns_str}\n\n"
        f"Mathematically Calculated KPIs & Metrics (Use these exact values in your analysis & report):\n"
        f"{kpis_text}\n\n"
        f"Data Column Types:\n{info_str}\n"
        f"Statistical Summary:\n{desc_str}\n\n"
        f"First 30 rows of the dataset:\n{sample_str}"
    )

    
    narrative_prompt = (
        "You are an elite corporate advisor, CFO business analyst, and expert data scientist.\n"
        "Your task is to analyze the following dataset that has been loaded from a file:\n\n"
        f"{dataset_context}\n\n"
        "Analyze this data and provide a highly engaging, professional, and plain-English Executive Narrative summarizing the key findings.\n"
        "Since the dataset type can vary (it could be sales, inventory, employee directory, tasks, marketing metrics, etc.), detect the type of data and adapt your analysis to it.\n\n"
        "CRITICAL INSTRUCTIONS FOR METRICS & NUMBERS:\n"
        "1. Report all currency values, revenue, sales, and profit figures EXACTLY as they are formatted in the 'Mathematically Calculated KPIs & Metrics' section.\n"
        "2. If no currency symbol (like '$' or '₹') is present in the calculated metrics, DO NOT add one. For example, if a number is given as '25,043.05', do not write it as '$25,043.05' or '₹25,043.05'. Just write '25,043.05'.\n"
        "3. Do not convert the numbers to abbreviations like '25k' or '25.04 thousand' or 'million'. Keep the precise numbers (like '25,043.05' or '1,250,000.00') exactly as provided in the metrics list.\n\n"
        "Format the response precisely with exactly 3 bullet points:\n"
        "1. **Data Overview & Health**: (A simplified takeaway of what this data represents, its size, completeness, and overall health/relevance)\n"
        "2. **Key Metrics & Highlights**: (Calculate and highlight the most critical metrics, top performers, averages, counts, or standout trends directly from the data)\n"
        "3. **Strategic Action Plan**: (Two high-value, specific, and actionable strategic recommendations based on these findings to improve efficiency, margins, or operations)\n"
        "Ensure the tone is friendly, elite, and completely free of raw technical jargon or mathematical formulas, so a non-technical end user understands it instantly."
    )
    
    st.session_state.ai_insights_narrative = ask_llm(narrative_prompt)

def attach_dataframe(file_obj, filename):
    st.session_state.cb_df = load_dataframe_from_file(file_obj, filename)
    st.session_state.cb_df_name = filename
    st.session_state.selected_nav = "Data Analysis"
    save_attachment_record(filename)
    with st.spinner("AI Business Advisor is reviewing your financials..."):
        auto_generate_insights()

def extract_pptx_text(file_obj):
    from pptx import Presentation

    presentation = Presentation(file_obj)
    slide_text = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slide_text.append(f"Slide {index}:\n" + "\n".join(parts))
    return "\n\n".join(slide_text)

def attach_file_to_chat_context(file_obj, filename, mime_type=""):
    name = filename.lower()
    st.session_state.cb_file_text = ""
    st.session_state.cb_img_bytes = None
    st.session_state.cb_img_mime = None
    st.session_state.cb_df = None
    st.session_state.cb_df_name = ""

    if name.endswith(".pdf"):
        raw = file_obj.read()
        from pdfminer3.layout import LAParams
        from pdfminer3.pdfpage import PDFPage
        from pdfminer3.pdfinterp import PDFResourceManager, PDFPageInterpreter
        from pdfminer3.converter import TextConverter

        rm = PDFResourceManager()
        sio = io.StringIO()
        cv = TextConverter(rm, sio, laparams=LAParams())
        pi = PDFPageInterpreter(rm, cv)
        for pg in PDFPage.get_pages(io.BytesIO(raw), caching=True, check_extractable=True):
            pi.process_page(pg)
        st.session_state.cb_file_text = sio.getvalue()
        cv.close()
        sio.close()
    elif name.endswith(".txt"):
        st.session_state.cb_file_text = file_obj.read().decode("utf-8", errors="ignore")
    elif name.endswith(".doc") or name.endswith(".docx"):
        import docx2txt
        st.session_state.cb_file_text = docx2txt.process(file_obj)
    elif name.endswith(".pptx"):
        st.session_state.cb_file_text = extract_pptx_text(file_obj)
    elif name.endswith((".csv", ".xlsx", ".xlsm", ".xls")):
        st.session_state.cb_df = load_dataframe_from_file(file_obj, filename)
        st.session_state.cb_df_name = filename
        with st.spinner("AI Business Advisor is reviewing your financials..."):
            auto_generate_insights()
    elif name.endswith((".png", ".jpg", ".jpeg", ".webp")):
        st.session_state.cb_img_bytes = file_obj.read()
        st.session_state.cb_img_mime = mime_type or "image/jpeg"
    else:
        raise ValueError("Unsupported file type. Use PDF, TXT, Word, PowerPoint, image, CSV, or Excel files.")

    save_attachment_record(filename, mime_type, source="google_drive" if filename.startswith("google_") else "chat_upload")
    
    # Ensure a conversation exists and log alternating turns for document import
    ensure_active_conversation()
    file_kind = get_file_type(filename)
    user_msg = f"Imported from Google Drive: {filename}"
    if file_kind == "document":
        assistant_msg = f"📄 **{filename}** has been successfully imported from Google Drive and attached to this conversation. You can now ask questions about this document!"
    elif file_kind == "dataset":
        assistant_msg = f"📊 **{filename}** has been successfully imported from Google Drive and loaded into the Data Analysis Workspace. What would you like to analyze?"
    elif file_kind == "image":
        assistant_msg = f"🖼️ **{filename}** has been successfully imported from Google Drive and attached. Ask me any questions about this image!"
    else:
        assistant_msg = f"📎 **{filename}** has been successfully imported from Google Drive and attached."
        
    st.session_state.cb_messages.append({
        "role": "user",
        "content": user_msg,
        "snippet": filename
    })
    st.session_state.cb_messages.append({
        "role": "assistant",
        "content": assistant_msg
    })
    if db_enabled():
        conversation_id = st.session_state.active_conversation_id
        db_action(save_message, int(conversation_id), "user", user_msg, filename, None, None)
        db_action(save_message, int(conversation_id), "assistant", assistant_msg, None, None, None)
    sync_active_conversation()
    st.session_state.selected_nav = "Chat"

def extract_drive_file_id(value):
    value = value.strip()
    patterns = [
        r"/d/([a-zA-Z0-9_-]+)",
        r"[?&]id=([a-zA-Z0-9_-]+)",
        r"^([a-zA-Z0-9_-]{20,})$",
    ]
    for pattern in patterns:
        match = re.search(pattern, value)
        if match:
            return match.group(1)
    return ""

def extract_filename_from_headers(headers, default_filename):
    cd = headers.get("Content-Disposition", "")
    if not cd:
        return default_filename
    
    # Try RFC 5987 UTF-8 encoding first
    rfc_match = re.search(r"filename\*=UTF-8''([^;\n]+)", cd, re.IGNORECASE)
    if rfc_match:
        try:
            return urllib.parse.unquote(rfc_match.group(1))
        except Exception:
            pass
            
    # Fallback to standard filename="filename"
    std_match = re.search(r'filename="?([^";\n]+)"?', cd, re.IGNORECASE)
    if std_match:
        return std_match.group(1)
        
    return default_filename

def download_drive_share_link(link_or_id, file_kind):
    file_id = extract_drive_file_id(link_or_id)
    if not file_id:
        raise ValueError("Paste a valid Google Drive share link or file ID.")

    mime_type = ""
    if file_kind == "Google Sheet":
        url = f"https://docs.google.com/spreadsheets/d/{file_id}/export?format=xlsx"
        filename = "google_sheet.xlsx"
        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif file_kind == "Google Doc":
        url = f"https://docs.google.com/document/d/{file_id}/export?format=docx"
        filename = "google_doc.docx"
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif file_kind == "Google Slides":
        url = f"https://docs.google.com/presentation/d/{file_id}/export/pptx"
        filename = "google_slides.pptx"
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif file_kind == "XLSX":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.xlsx"
        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif file_kind == "XLS":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.xls"
        mime_type = "application/vnd.ms-excel"
    elif file_kind == "CSV":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.csv"
        mime_type = "text/csv"
    elif file_kind == "PDF":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.pdf"
        mime_type = "application/pdf"
    elif file_kind == "Text":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.txt"
        mime_type = "text/plain"
    elif file_kind == "Word document":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.docx"
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif file_kind == "PPTX":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.pptx"
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif file_kind == "PNG":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_image.png"
        mime_type = "image/png"
    elif file_kind == "JPG/JPEG":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_image.jpg"
        mime_type = "image/jpeg"
    elif file_kind == "WEBP":
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_image.webp"
        mime_type = "image/webp"
    else:
        url = f"https://drive.google.com/uc?export=download&id={file_id}"
        filename = "google_drive_file.xlsx"
        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

    request = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(request, timeout=60) as response:
        data = response.read()
        content_type = response.headers.get("Content-Type", "")
        # Extract original filename from Content-Disposition header if available
        filename = extract_filename_from_headers(response.headers, filename)

    if b"Google Drive - Virus scan warning" in data or "text/html" in content_type:
        raise ValueError("Could not download this file. Make sure the Drive link is shared with access set to Anyone with the link.")

    return io.BytesIO(data), filename, mime_type

def render_data_analysis():
    import pandas as pd
    col1, col2 = st.columns([9, 1])
    with col1:
        st.subheader("Data Analysis Workspace")
    with col2:
        if st.button("❌", key="close_data_ws", use_container_width=True):
            st.session_state.selected_nav = "Chat"
            st.rerun()
            
    if st.session_state.cb_df is None:
        st.info("No dataset is attached yet. Upload a CSV or Excel file in Chat or from Google Drive.")
        return

    df_loaded = st.session_state.cb_df

    # Basic state initialization
    if "chart_interpretation" not in st.session_state:
        st.session_state.chart_interpretation = ""
    if "prev_chart_config" not in st.session_state:
        st.session_state.prev_chart_config = {}

    # Advanced Heuristic calculations
    mapping = auto_detect_columns(df_loaded)
    sales_col = mapping["sales"]
    profit_col = mapping["profit"]
    cust_col = mapping["customer"]
    date_col = mapping["date"]

    active_cols = [c for c in [sales_col, profit_col] if c]
    currency_symbol = detect_currency_symbol(df_loaded, active_cols)
    total_sales = 0.0
    total_profit = 0.0
    margin = 0.0
    transactions = len(df_loaded)
    avg_order = 0.0
    avg_profit = 0.0
    max_order_val = 0.0
    unique_customers = 0
    top_customer_name = "N/A"
    top_customer_sales = 0.0

    if sales_col:
        sales_series = clean_numeric_column(df_loaded, sales_col)
        total_sales = float(sales_series.sum())
        avg_order = float(sales_series.mean())
        max_order_val = float(sales_series.max())
        
    if profit_col:
        profit_series = clean_numeric_column(df_loaded, profit_col)
        total_profit = float(profit_series.sum())
        if transactions > 0:
            avg_profit = total_profit / transactions
            
    if total_sales > 0.0:
        margin = (total_profit / total_sales) * 100

    if cust_col:
        unique_customers = int(df_loaded[cust_col].nunique())
        if sales_col:
            cust_col_data = df_loaded[cust_col].astype(str)
            sales_col_data = clean_numeric_column(df_loaded, sales_col)
            cust_df = pd.DataFrame({cust_col: cust_col_data, sales_col: sales_col_data})
            cust_totals = cust_df.groupby(cust_col)[sales_col].sum()
            if not cust_totals.empty:
                top_customer_name = str(cust_totals.idxmax())
                top_customer_sales = float(cust_totals.max())


    st.markdown('<h3 style="margin-top:0px; color:#da7756;">Business Performance & KPI Summary</h3>', unsafe_allow_html=True)
    st.caption("Auto-calculated critical metrics and plain-English intelligence from your dataset.")
    
    # Display styled KPI summary cards dynamically
    st.markdown('<div style="font-weight:700; color:#da7756; margin-bottom:16px; font-size:1.1rem; border-bottom: 1px solid #e5e3d9; padding-bottom: 6px;">📊 Key Performance Indicators (KPIs)</div>', unsafe_allow_html=True)
    
    metrics_to_show = []
    metrics_to_show.append(("Total Orders (Transactions)", f"{transactions:,}", "Total number of transaction records"))
    
    if sales_col:
        metrics_to_show.append(("Total Revenue (Sales)", f"{currency_symbol}{total_sales:,.2f}", "Sum of transaction sales. Represents top-line growth."))
        metrics_to_show.append(("Average Order Value (AOV)", f"{currency_symbol}{avg_order:,.2f}", "Average spend per transaction. Indicates pricing strength."))
        metrics_to_show.append(("Max Single Order Value", f"{currency_symbol}{max_order_val:,.2f}", "Largest single sales order recorded."))
        
    if profit_col:
        metrics_to_show.append(("Net Profit", f"{currency_symbol}{total_profit:,.2f}", "Net profit earned. Represents bottom-line earnings."))
        metrics_to_show.append(("Average Profit per Order", f"{currency_symbol}{avg_profit:,.2f}", "Average net profit margins generated per transaction."))
        
    if sales_col and profit_col:
        metrics_to_show.append(("Net Profit Margin", f"{margin:.2f}%", "Net Profit percentage of Revenue. Shows overall cost efficiency."))
        
    if cust_col:
        metrics_to_show.append(("Total Unique Customers", f"{unique_customers:,}", "Number of distinct buyers. Tracks customer acquisition."))

    # Render metrics dynamically in rows of 3 columns
    for idx in range(0, len(metrics_to_show), 3):
        cols = st.columns(min(3, len(metrics_to_show) - idx))
        for i, col_widget in enumerate(cols):
            label, value, help_text = metrics_to_show[idx + i]
            col_widget.metric(label, value, help=help_text)

    # AI Executive Business Narrative Section
    st.markdown("---")
    
    if not st.session_state.get("ai_insights_narrative"):
        with st.spinner("AI Business Advisor is reviewing your financials..."):
            auto_generate_insights()

    narrative = st.session_state.get("ai_insights_narrative", "")
    if narrative:
        st.markdown(
            f"""
            <div class="ai-insights-card">
                <div class="ai-insights-title">💡 AI Executive Business Insights</div>
            """,
            unsafe_allow_html=True
        )
        st.markdown(narrative)
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.warning("Could not generate AI insights. Please check your API keys or configurations.")



def find_in_memory_content(item):
    conv_id = item["conversation_id"]
    ftype = item["file_type"]
    for conv in st.session_state.conversations:
        if conv["id"] == conv_id:
            if ftype == "dataset":
                df = conv.get("cb_df")
                data_json = df.to_json(orient="records") if df is not None else None
                return {"data_json": data_json}
            elif ftype == "image":
                return {"image_data": conv.get("cb_img_bytes")}
            elif ftype == "document":
                return {"text_content": conv.get("cb_file_text")}
    return None
def restore_file_to_chat(item):
    import pandas as pd
    ftype = item["file_type"]
    filename = item["filename"]
    
    if db_enabled():
        content_row = db_action(load_upload_content, int(item["id"]))
    else:
        content_row = find_in_memory_content(item)
        
    if not content_row:
        st.error("Failed to retrieve content for restoring.")
        return
        
    # Start a brand new chat session to restore the file into
    start_new_chat()
        
    st.session_state.cb_df = None
    st.session_state.cb_df_name = ""
    st.session_state.cb_file_text = ""
    st.session_state.cb_img_bytes = None
    st.session_state.cb_img_mime = None
    st.session_state.ai_insights_narrative = ""

    if ftype == "dataset":
        if content_row.get("data_json"):
            import json
            data = json.loads(content_row["data_json"])
            st.session_state.cb_df = pd.DataFrame(data)
            st.session_state.cb_df_name = filename
            st.session_state.selected_nav = "Chat"
            auto_generate_insights()
    elif ftype == "image":
        st.session_state.cb_img_bytes = content_row.get("image_data")
        st.session_state.cb_img_mime = item.get("mime_type") or "image/jpeg"
        st.session_state.selected_nav = "Chat"
    elif ftype == "document":
        st.session_state.cb_file_text = content_row.get("text_content") or ""
        st.session_state.selected_nav = "Chat"

    # Add message and sync conversation so it shows up in Recents list
    ensure_active_conversation()
    file_kind = get_file_type(filename)
    user_msg = f"Restored file: {filename}"
    if file_kind == "document":
        assistant_msg = f"📄 **{filename}** has been successfully restored to this conversation. You can now ask questions about this document!"
    elif file_kind == "dataset":
        assistant_msg = f"📊 **{filename}** has been successfully restored and loaded into the Data Analysis Workspace."
    elif file_kind == "image":
        assistant_msg = f"🖼️ **{filename}** has been successfully restored to this conversation. Ask me any questions about this image!"
    else:
        assistant_msg = f"📎 **{filename}** has been successfully restored to this conversation."

    st.session_state.cb_messages.append({
        "role": "user",
        "content": user_msg,
        "snippet": filename
    })
    st.session_state.cb_messages.append({
        "role": "assistant",
        "content": assistant_msg
    })
    if db_enabled():
        conversation_id = st.session_state.active_conversation_id
        db_action(
            save_message,
            int(conversation_id),
            "user",
            user_msg,
            filename,
            None,
            None
        )
        db_action(
            save_message,
            int(conversation_id),
            "assistant",
            assistant_msg,
            None,
            None,
            None
        )
        # Save attachment record to associate it with the new conversation
        save_attachment_record(filename, item.get("mime_type") or "", source=item.get("source") or "chat_upload")
    sync_active_conversation()

def delete_file_record(item):
    item_id = item["id"]
    if db_enabled():
        db_action(delete_upload, int(item_id))
    else:
        conv_id = item["conversation_id"]
        ftype = item["file_type"]
        for conv in st.session_state.conversations:
            if conv["id"] == conv_id:
                if ftype == "dataset":
                    conv["cb_df"] = None
                    conv["cb_df_name"] = ""
                elif ftype == "image":
                    conv["cb_img_bytes"] = None
                elif ftype == "document":
                    conv["cb_file_text"] = ""
                break
                
    if st.session_state.get("preview_file_id") == item_id:
        st.session_state.preview_file_id = None


def render_document_library():
    import pandas as pd
    col1, col2 = st.columns([9, 1])
    with col1:
        st.subheader("Document & File Library")
    with col2:
        if st.button("❌", key="close_doc_lib", use_container_width=True):
            st.session_state.selected_nav = "Chat"
            st.rerun()
    st.caption("Manage your historically uploaded documents, images, and datasets. Click 'Preview' to view file contents, 'Download' to download a file, 'Restore' to load a file into your active chat session, or 'Delete' to remove it.")

    uploads = []
    if db_enabled():
        uploads = cached_load_user_uploads(st.session_state.db_user_id) or []
    else:
        for conversation in st.session_state.conversations:
            conv_id = conversation["id"]
            conv_title = conversation.get("title", "Conversation")
            if conversation.get("cb_df") is not None:
                uploads.append({
                    "id": f"mem_df_{conv_id}",
                    "conversation_id": conv_id,
                    "conv_title": conv_title,
                    "filename": conversation.get("cb_df_name", "dataset.csv"),
                    "file_type": "dataset",
                    "mime_type": "text/csv",
                    "source": "chat_upload",
                    "uploaded_at": conversation.get("updated_at", ""),
                    "rows_count": len(conversation["cb_df"]),
                    "columns_count": len(conversation["cb_df"].columns),
                })
            if conversation.get("cb_img_bytes") is not None:
                uploads.append({
                    "id": f"mem_img_{conv_id}",
                    "conversation_id": conv_id,
                    "conv_title": conv_title,
                    "filename": "image.jpg",
                    "file_type": "image",
                    "mime_type": conversation.get("cb_img_mime", "image/jpeg"),
                    "source": "chat_upload",
                    "uploaded_at": conversation.get("updated_at", ""),
                })
            if conversation.get("cb_file_text"):
                uploads.append({
                    "id": f"mem_doc_{conv_id}",
                    "conversation_id": conv_id,
                    "conv_title": conv_title,
                    "filename": "document.txt",
                    "file_type": "document",
                    "mime_type": "text/plain",
                    "source": "chat_upload",
                    "uploaded_at": conversation.get("updated_at", ""),
                })

    if not uploads:
        st.info("No files uploaded yet. Files you upload in Chat or import from Google Drive will appear here.")
        return

    tab_all, tab_data, tab_docs, tab_imgs = st.tabs([
        "📁 All Files", "📊 Datasets", "📄 Documents", "🖼️ Images"
    ])

    def render_file_list(filtered_uploads, category_name):
        if not filtered_uploads:
            st.info(f"No {category_name} found.")
            return

        for idx, item in enumerate(filtered_uploads):
            col_icon, col_details, col_actions = st.columns([0.6, 4.4, 5.0])
            
            with col_icon:
                ftype = item["file_type"]
                if ftype == "dataset":
                    st.markdown('<div style="text-align: center; margin-top: 10px;"><i class="fa-solid fa-table" style="color: #da7756; font-size: 1.6rem;"></i></div>', unsafe_allow_html=True)
                elif ftype == "image":
                    st.markdown('<div style="text-align: center; margin-top: 10px;"><i class="fa-solid fa-image" style="color: #ec4899; font-size: 1.6rem; filter: drop-shadow(0 2px 4px rgba(236, 72, 153, 0.3));"></i></div>', unsafe_allow_html=True)
                else:
                    st.markdown('<div style="text-align: center; margin-top: 10px;"><i class="fa-solid fa-file-lines" style="color: #10b981; font-size: 1.6rem; filter: drop-shadow(0 2px 4px rgba(16, 185, 129, 0.3));"></i></div>', unsafe_allow_html=True)

            with col_details:
                filename = item["filename"]
                uploaded_at = item["uploaded_at"]
                
                import datetime
                import os
                
                local_uploaded_at = uploaded_at
                if isinstance(uploaded_at, str):
                    try:
                        # Clean string and convert to datetime object
                        clean_str = uploaded_at.replace("T", " ").replace("Z", "").split(".")[0]
                        local_uploaded_at = datetime.datetime.strptime(clean_str, "%Y-%m-%d %H:%M:%S")
                    except Exception:
                        pass
                
                if isinstance(local_uploaded_at, datetime.datetime):
                    if local_uploaded_at.tzinfo is None:
                        local_uploaded_at = local_uploaded_at.replace(tzinfo=datetime.timezone.utc)
                    
                    if os.name == 'nt':
                        local_uploaded_at = local_uploaded_at.astimezone()
                    else:
                        ist_tz = datetime.timezone(datetime.timedelta(hours=5, minutes=30))
                        local_uploaded_at = local_uploaded_at.astimezone(ist_tz)
                        
                    date_str = local_uploaded_at.strftime("%Y-%m-%d %H:%M:%S")
                else:
                    date_str = str(uploaded_at)

                size_info = ""

                source_label = "Google Drive" if item["source"] == "google_drive" else "Chat Upload"
                
                st.markdown(
                    f'<div style="margin-top: 2px;">'
                    f'<span style="font-weight: 700; color: #1c1b1a; font-size: 1.05rem;">{filename}</span>'
                    f'{size_info}<br>'
                    f'<span style="color: #6b685c; font-size: 0.82rem;">Uploaded: {date_str} &nbsp;|&nbsp; Source: <strong>{source_label}</strong></span>'
                    f'</div>', 
                    unsafe_allow_html=True
                )

            with col_actions:
                col_btn1, col_btn2, col_btn3, col_btn4 = st.columns(4)
                item_id = item["id"]
                
                with col_btn1:
                    if st.button("Preview", key=f"preview_btn_{category_name}_{item_id}_{idx}", use_container_width=True):
                        st.session_state.preview_file_id = item_id
                        st.rerun()
                
                with col_btn2:
                    # Fetch file content for direct download from the list
                    if db_enabled():
                        content_row = db_action(load_upload_content, int(item_id))
                    else:
                        content_row = find_in_memory_content(item)
                    
                    download_data = None
                    download_filename = item["filename"]
                    download_mime = item.get("mime_type") or "application/octet-stream"
                    
                    if content_row:
                        if ftype == "dataset" and content_row.get("data_json"):
                            try:
                                import json
                                data = json.loads(content_row["data_json"])
                                df = pd.DataFrame(data)
                                download_data = df.to_csv(index=False).encode("utf-8")
                                download_mime = "text/csv"
                                if not download_filename.lower().endswith(".csv"):
                                    download_filename = os.path.splitext(download_filename)[0] + ".csv"
                            except Exception:
                                download_data = content_row["data_json"].encode("utf-8")
                                download_mime = "application/json"
                        elif ftype == "image" and content_row.get("image_data"):
                            download_data = content_row["image_data"]
                        elif ftype == "document" and content_row.get("text_content"):
                            download_data = content_row["text_content"].encode("utf-8")
                            download_mime = "text/plain"
                            if not download_filename.lower().endswith((".txt", ".pdf", ".docx", ".doc", ".pptx")):
                                download_filename = os.path.splitext(download_filename)[0] + "_extracted.txt"

                    if download_data is not None:
                        st.download_button(
                            label="Download",
                            data=download_data,
                            file_name=download_filename,
                            mime=download_mime,
                            key=f"download_btn_{category_name}_{item_id}_{idx}",
                            use_container_width=True
                        )
                    else:
                        st.button("Download", key=f"download_disabled_{category_name}_{item_id}_{idx}", disabled=True, use_container_width=True)

                with col_btn3:
                    if st.button("Restore", key=f"restore_btn_{category_name}_{item_id}_{idx}", use_container_width=True):
                        restore_file_to_chat(item)
                        st.success("File restored in a new chat session!")
                        st.rerun()

                with col_btn4:
                    if st.button("Delete", key=f"delete_btn_{category_name}_{item_id}_{idx}", type="secondary", use_container_width=True):
                        delete_file_record(item)
                        st.success("File deleted successfully!")
                        st.rerun()

            st.markdown('<hr style="margin: 8px 0; border: 0; border-top: 1px solid rgba(255, 255, 255, 0.05);">', unsafe_allow_html=True)

    datasets = [u for u in uploads if u["file_type"] == "dataset"]
    documents = [u for u in uploads if u["file_type"] == "document"]
    images = [u for u in uploads if u["file_type"] == "image"]

    with tab_all:
        render_file_list(uploads, "files")
    with tab_data:
        render_file_list(datasets, "datasets")
    with tab_docs:
        render_file_list(documents, "documents")
    with tab_imgs:
        render_file_list(images, "images")

    preview_id = st.session_state.get("preview_file_id")
    if preview_id:
        st.markdown('<div id="preview-anchor"></div>', unsafe_allow_html=True)
        st.markdown("---")
        item = None
        for u in uploads:
            if u["id"] == preview_id:
                item = u
                break
        
        if item:
            st.markdown(f'<h3 style="color: #da7756; margin-bottom: 12px;"><i class="fa-solid fa-magnifying-glass"></i> File Preview: {item["filename"]}</h3>', unsafe_allow_html=True)
            
            with st.container(border=True):
                ftype = item["file_type"]
                
                if db_enabled():
                    content_row = db_action(load_upload_content, int(item["id"]))
                else:
                    content_row = find_in_memory_content(item)

                if not content_row:
                    st.error("Could not load preview content.")
                else:
                    if ftype == "dataset":
                        if content_row.get("data_json"):
                            try:
                                import json
                                data = json.loads(content_row["data_json"])
                                df = pd.DataFrame(data)
                                st.markdown("**Dataset Preview**")
                                if len(df) > 1000:
                                    st.dataframe(df.head(1000), use_container_width=True)
                                    st.caption(f"⚠️ Showing first 1,000 rows out of {len(df):,} total rows to prevent frontend lag.")
                                else:
                                    st.dataframe(df, use_container_width=True)
                            except Exception as ex:
                                st.error(f"Error parsing dataset preview: {ex}")
                        else:
                            st.warning("No tabular preview data found for this dataset.")
                            
                    elif ftype == "image":
                        image_data = content_row.get("image_data")
                        if image_data:
                            try:
                                st.image(io.BytesIO(image_data), use_container_width=True, caption=item["filename"])
                            except Exception as ex:
                                st.error(f"Error rendering image: {ex}")
                        else:
                            st.warning("No image data found to preview.")
                            
                    elif ftype == "document":
                        text_content = content_row.get("text_content")
                        if text_content:
                            st.text_area("Extracted Document Text", text_content, height=380)
                        else:
                            st.warning("This document is empty or has no extracted text.")

                    # Download button logic inside the container
                    download_data = None
                    download_filename = item["filename"]
                    download_mime = item.get("mime_type") or "application/octet-stream"
                    
                    if ftype == "dataset" and content_row.get("data_json"):
                        try:
                            import json
                            data = json.loads(content_row["data_json"])
                            df = pd.DataFrame(data)
                            download_data = df.to_csv(index=False).encode("utf-8")
                            download_mime = "text/csv"
                            if not download_filename.lower().endswith(".csv"):
                                download_filename = os.path.splitext(download_filename)[0] + ".csv"
                        except Exception:
                            download_data = content_row["data_json"].encode("utf-8")
                            download_mime = "application/json"
                    elif ftype == "image" and content_row.get("image_data"):
                        download_data = content_row["image_data"]
                    elif ftype == "document" and content_row.get("text_content"):
                        download_data = content_row["text_content"].encode("utf-8")
                        download_mime = "text/plain"
                        if not download_filename.lower().endswith((".txt", ".pdf", ".docx", ".doc", ".pptx")):
                            download_filename = os.path.splitext(download_filename)[0] + "_extracted.txt"

                    if download_data is not None:
                        st.markdown("<div style='margin-top: 15px;'></div>", unsafe_allow_html=True)
                        st.download_button(
                            label=f"📥 Download {download_filename}",
                            data=download_data,
                            file_name=download_filename,
                            mime=download_mime,
                            key=f"download_file_{preview_id}",
                            use_container_width=True
                        )
                            
            if st.button("Close Preview", key="close_preview_btn"):
                st.session_state.preview_file_id = None
                st.rerun()

            # Auto scroll to preview section using inline JS bypass
            st.markdown(
                """
                <img src="x" style="opacity: 0; width: 0; height: 0; position: absolute;" onerror="
                    this.style.display='none';
                    var element = document.getElementById('preview-anchor');
                    if (element) {
                        element.scrollIntoView({ behavior: 'smooth', block: 'start' });
                    }
                " />
                """,
                unsafe_allow_html=True
            )

def render_webhooks_integration():
    col1, col2 = st.columns([9, 1])
    with col1:
        st.subheader("Webhook Integrations")
    with col2:
        if st.button("❌", key="close_webhooks_ws", use_container_width=True):
            st.session_state.selected_nav = "Chat"
            st.rerun()
            
    st.caption("Setup and monitor Telegram and WhatsApp Webhooks to interact with your AI assistant on messaging platforms.")

    telegram_token = get_secret("TELEGRAM_BOT_TOKEN", "")
    whatsapp_token = get_secret("WHATSAPP_ACCESS_TOKEN", "")
    whatsapp_phone_id = get_secret("WHATSAPP_PHONE_NUMBER_ID", "")
    whatsapp_verify_token = get_secret("WHATSAPP_VERIFY_TOKEN", "")

    # Status indicators in columns
    c1, c2 = st.columns(2)
    
    with c1:
        st.markdown("### 🤖 Telegram Bot")
        if telegram_token and telegram_token != "your-telegram-bot-token":
            st.success("🟢 Connected")
            st.markdown(f"**Bot Token:** `{telegram_token[:6]}...{telegram_token[-6:]}`")
        else:
            st.warning("🟡 Not Configured")
            st.info("To configure Telegram, add `TELEGRAM_BOT_TOKEN` in your secrets/environment variables.")

    with c2:
        st.markdown("### 💬 WhatsApp Business Cloud")
        if (whatsapp_token and whatsapp_token != "your-whatsapp-access-token" and
            whatsapp_phone_id and whatsapp_phone_id != "your-whatsapp-phone-number-id" and
            whatsapp_verify_token and whatsapp_verify_token != "your-whatsapp-webhook-verify-token"):
            st.success("🟢 Connected")
            st.markdown(f"**Phone Number ID:** `{whatsapp_phone_id}`")
            st.markdown(f"**Verify Token:** `{whatsapp_verify_token}`")
        else:
            st.warning("🟡 Not Configured")
            st.info("To configure WhatsApp, add `WHATSAPP_ACCESS_TOKEN`, `WHATSAPP_PHONE_NUMBER_ID`, and `WHATSAPP_VERIFY_TOKEN` in your secrets/environment variables.")

    st.markdown("---")
    
    tab_urls, tab_tg_guide, tab_wa_guide = st.tabs([
        "🔌 Callback URLs", "🤖 Telegram Setup Guide", "💬 WhatsApp Setup Guide"
    ])

    with tab_urls:
        st.markdown("### Webhook Callback Endpoints")
        st.markdown(
            "Use a public tunneling tool like **ngrok** to forward requests to your local FastAPI backend port `8000` "
            "if you are running the project locally."
        )
        
        # Display the local endpoints
        st.info("FastAPI Backend webhook routes:")
        st.markdown(
            "- **Telegram Callback URL:**\n"
            "  `http://127.0.0.1:8000/webhook/telegram` (POST)\n"
            "- **WhatsApp Callback URL:**\n"
            "  `http://127.0.0.1:8000/webhook/whatsapp` (GET / POST)\n"
        )
        
        st.markdown(
            "💡 **Note:** Replace `http://127.0.0.1:8000` with your public ngrok URL (e.g. `https://xxxx-xx-xx.ngrok-free.app`) "
            "when registering webhooks on Meta / Telegram developer portals."
        )

    with tab_tg_guide:
        st.markdown("### How to Setup Telegram Bot Webhook")
        st.markdown(
            """
            1. Find **@BotFather** on Telegram and send `/newbot` command to create a bot.
            2. Copy the **HTTP API Token** and save it as `TELEGRAM_BOT_TOKEN` in your `.streamlit/secrets.toml`.
            3. Open your browser or run a curl command to set the webhook callback:
               ```bash
               curl -X POST "https://api.telegram.org/bot<YOUR_TELEGRAM_BOT_TOKEN>/setWebhook?url=<YOUR_PUBLIC_TUNNEL_URL>/webhook/telegram"
               ```
            4. Once registered successfully, any message sent to your bot will be answered by this assistant.
            """
        )

    with tab_wa_guide:
        st.markdown("### How to Setup WhatsApp Cloud Webhook")
        st.markdown(
            """
            1. Go to **[Meta Developers Portal](https://developers.facebook.com/)** and create a **Business App**.
            2. Add **WhatsApp** product to your app.
            3. Under WhatsApp > **API Setup**, find your **Phone Number ID** and **Temporary Access Token** (or create a permanent system user token).
            4. Add these values along with a custom **Verify Token** to your `.streamlit/secrets.toml`.
            5. In the Meta App Dashboard, go to **WhatsApp > Configuration**:
               - Click **Edit** under Webhooks.
               - Enter **Callback URL**: `<YOUR_PUBLIC_TUNNEL_URL>/webhook/whatsapp`
               - Enter **Verify Token**: (the verify token value you set in secrets)
               - Click **Verify and save**.
            6. Subscribe to **messages** webhook fields under WhatsApp Webhook Fields table.
            """
        )

def render_google_drive_upload():
    col1, col2 = st.columns([9, 1])
    with col1:
        st.subheader("Upload From Google Drive")
    with col2:
        if st.button("❌", key="close_gdrive", use_container_width=True):
            st.session_state.selected_nav = "Chat"
            st.rerun()
    st.caption("Browse and select files directly from your Google Drive using your Google credentials.")

    if not GOOGLE_CLIENT_ID:
        st.warning("⚠️ Google Client ID is not configured. Direct browsing is disabled. Please configure GOOGLE_CLIENT_ID in your environment variables.")
    else:
        # Load developer key (default fallback or environment variable)
        developer_key = get_secret("GOOGLE_DEVELOPER_KEY", "")
        app_id = GOOGLE_CLIENT_ID.split("-")[0] if "-" in GOOGLE_CLIENT_ID else ""

        # Display the custom Google Drive Picker component
        if google_drive_picker:
            try:
                # This component renders the "Browse Google Drive" button
                picked_file = google_drive_picker(
                    clientId=GOOGLE_CLIENT_ID,
                    developerKey=developer_key,
                    appId=app_id,
                    key="google_drive_picker_instance"
                )
                
                if picked_file:
                    file_id = picked_file.get("id")
                    filename = picked_file.get("name", "Unnamed File")
                    mime_type = picked_file.get("mimeType", "")
                    token = picked_file.get("token")
                    
                    if file_id and token:
                        with st.spinner(f"📥 Downloading and importing '{filename}' from Google Drive..."):
                            try:
                                headers = {"Authorization": f"Bearer {token}"}
                                
                                # Handle Google Workspace files that must be exported
                                if mime_type == "application/vnd.google-apps.document":
                                    download_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                    actual_filename = filename if filename.lower().endswith(".docx") else filename + ".docx"
                                    actual_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                elif mime_type == "application/vnd.google-apps.spreadsheet":
                                    download_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                                    actual_filename = filename if filename.lower().endswith(".xlsx") else filename + ".xlsx"
                                    actual_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                                elif mime_type == "application/vnd.google-apps.presentation":
                                    download_url = f"https://www.googleapis.com/drive/v3/files/{file_id}/export?mimeType=application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                    actual_filename = filename if filename.lower().endswith(".pptx") else filename + ".pptx"
                                    actual_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                else:
                                    # Standard binary download for other files
                                    download_url = f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media"
                                    actual_filename = filename
                                    actual_mime = mime_type
                                
                                response = requests.get(download_url, headers=headers, timeout=60)
                                if response.status_code == 200:
                                    file_obj = io.BytesIO(response.content)
                                    attach_file_to_chat_context(file_obj, actual_filename, actual_mime)
                                    st.session_state.drive_import_notice = f"Imported {actual_filename} into the chatbot. Ask your next question in Chat."
                                    st.success(f"🎉 Successfully imported {actual_filename} into the chatbot!")
                                    st.rerun()
                                else:
                                    st.error(f"Google Drive API returned an error (HTTP {response.status_code}): {response.text}")
                            except Exception as ex:
                                st.error(f"Failed to process download from Google Drive: {ex}")
            except Exception as e:
                st.error(f"Error loading custom component: {e}")
        else:
            st.error("Google Drive Picker custom component is not available.")

    st.markdown("---")
    with st.expander("Paste Shared Google Drive Link"):
        st.markdown(
            """
            <div style="font-size: 0.85rem; color: #94a3b8; margin-bottom: 12px;">
                Use this fallback if you prefer to paste a public link. The file must be shared as 
                <strong>'Anyone with the link can view'</strong>.
            </div>
            """, 
            unsafe_allow_html=True
        )
        @safe_fragment
        def drive_link_paste_block():
            drive_link = st.text_input(
                "Google Drive link",
                placeholder="Paste a shared Google Drive link",
                key="drive_share_link",
            )
            file_kind = st.selectbox(
                "Drive file type",
                ["Google Sheet", "Google Doc", "Google Slides", "XLSX", "XLS", "CSV", "PDF", "Text", "Word document", "PPTX", "PNG", "JPG/JPEG", "WEBP"],
                key="drive_file_kind",
            )
            if st.button("Import via Link", use_container_width=True):
                try:
                    file_obj, filename, mime_type = download_drive_share_link(drive_link, file_kind)
                    attach_file_to_chat_context(file_obj, filename, mime_type)
                    st.session_state.drive_import_notice = f"Imported {filename} into the chatbot. Ask your next question in Chat."
                    st.success(f"Imported {filename} into the chatbot.")
                    st.rerun()
                except Exception as ex:
                    st.error(f"Google Drive import error: {ex}")
                    
        drive_link_paste_block()

def render_github_explorer():
    col1, col2 = st.columns([9, 1])
    with col1:
        st.subheader("GitHub Repository Explorer")
    with col2:
        if st.button("Close", key="close_github_explorer", use_container_width=True):
            st.session_state.selected_nav = "Chat"
            st.rerun()
            
    st.caption("Browse repository structure and select source code files to import as context for your AI Assistant.")
    
    github_creds = st.session_state.get("github_credentials")
    if not github_creds:
        st.warning("GitHub not connected. Please connect your account in the sidebar.")
        st.stop()
        
    access_token = github_creds.get("access_token")
    if not access_token:
        st.error("Invalid credentials. Try disconnecting and reconnecting in the sidebar.")
        st.stop()
        
    from backend.github_service import list_repositories, list_repo_contents, get_repo_file_content
    
    with st.spinner("Loading repositories..."):
        repos = list_repositories(access_token)
        
    if not repos:
        st.info("No repositories found or access token expired. Verify your permissions.")
        st.stop()
        
    repo_options = [r["full_name"] for r in repos]
    
    # Track current selected repo index
    default_idx = 0
    if st.session_state.github_explorer_repo in repo_options:
        default_idx = repo_options.index(st.session_state.github_explorer_repo)
        
    selected_repo = st.selectbox(
        "Select Repository",
        repo_options,
        index=default_idx,
        key="github_repo_select_widget"
    )
    
    if selected_repo != st.session_state.github_explorer_repo:
        st.session_state.github_explorer_repo = selected_repo
        st.session_state.github_explorer_path = ""
        st.rerun()
        
    repo_name = st.session_state.github_explorer_repo
    path = st.session_state.github_explorer_path
    
    # Breadcrumbs UI
    st.markdown("### Directory Path")
    parts = [p for p in path.split("/") if p]
    
    # Styled breadcrumbs
    bc_html = '<div style="display: flex; flex-wrap: wrap; align-items: center; gap: 8px; margin-bottom: 15px; font-family: monospace; font-size: 0.95rem;">'
    bc_html += '<span style="color: #6b685c;">Path:</span>'
    bc_html += '<span style="background: #f4f4f5; padding: 4px 8px; border-radius: 4px; color: #1c1b1a; font-weight: bold;">[root]</span>'
    for p in parts:
        bc_html += '<span style="color: #94a3b8;">/</span>'
        bc_html += f'<span style="background: #f4f4f5; padding: 4px 8px; border-radius: 4px; color: #1c1b1a;">{p}</span>'
    bc_html += '</div>'
    st.markdown(bc_html, unsafe_allow_html=True)
    
    # Navigation Buttons for Breadcrumbs
    bc_cols = st.columns(max(len(parts) + 1, 1))
    with bc_cols[0]:
        if st.button("📁 Root", key="bc_root_btn", use_container_width=True):
            st.session_state.github_explorer_path = ""
            st.rerun()
            
    running_path = ""
    for idx, p in enumerate(parts):
        running_path = f"{running_path}/{p}" if running_path else p
        with bc_cols[idx + 1]:
            # Create button closure to capture the target path correctly
            def make_bc_cb(target_path):
                return lambda: setattr(st.session_state, "github_explorer_path", target_path)
            st.button(p, key=f"bc_btn_{idx}", on_click=make_bc_cb(running_path), use_container_width=True)
            
    st.markdown("---")
    
    with st.spinner("Fetching contents..."):
        contents = list_repo_contents(repo_name, path, access_token)
        
    if not contents:
        st.info("This folder is empty or couldn't be loaded.")
        st.stop()
        
    # Sort folders first, then files
    folders = [item for item in contents if item.get("type") == "dir"]
    files = [item for item in contents if item.get("type") == "file"]
    
    # Back button if we are not at root
    if path:
        parent_path = "/".join(parts[:-1])
        if st.button("Go Up (Parent Directory)", key="go_up_dir_btn"):
            st.session_state.github_explorer_path = parent_path
            st.rerun()
            
    # Display folders
    if folders:
        st.markdown("#### Folders")
        # Layout in 3-column grid for folders
        folder_cols = st.columns(3)
        for idx, folder in enumerate(folders):
            col_idx = idx % 3
            folder_name = folder.get("name")
            folder_path = folder.get("path")
            with folder_cols[col_idx]:
                def make_folder_cb(target_path):
                    return lambda: setattr(st.session_state, "github_explorer_path", target_path)
                st.button(f"📁 {folder_name}", key=f"folder_{idx}", on_click=make_folder_cb(folder_path), use_container_width=True)
                
    # Display files
    if files:
        st.markdown("#### Files")
        for idx, file_item in enumerate(files):
            file_name = file_item.get("name")
            file_path = file_item.get("path")
            
            # Show file row
            with st.container():
                fcol1, fcol2, fcol3 = st.columns([6, 2, 2])
                with fcol1:
                    st.markdown(f"📄 **{file_name}**")
                    st.caption(f"Path: {file_path}")
                with fcol2:
                    if st.button("Preview", key=f"preview_file_{idx}", use_container_width=True):
                        st.session_state[f"preview_{file_path}"] = True
                with fcol3:
                    if st.button("Attach to Chat", key=f"attach_file_{idx}", use_container_width=True):
                        with st.spinner(f"Downloading {file_name}..."):
                            file_content = get_repo_file_content(repo_name, file_path, access_token)
                            if file_content is not None:
                                import_github_file_to_chat_context(file_content, file_name, repo_name)
                                st.session_state.selected_nav = "Chat"
                                st.success(f"Attached {file_name} to Chat!")
                                time.sleep(1)
                                st.rerun()
                            else:
                                st.error(f"Failed to fetch content for {file_name}. It may be a binary file or too large.")
                                
            # Preview expander if active
            if st.session_state.get(f"preview_{file_path}", False):
                with st.spinner("Loading preview..."):
                    file_content = get_repo_file_content(repo_name, file_path, access_token)
                if file_content is not None:
                    with st.expander(f"Preview: {file_name}", expanded=True):
                        # Detect syntax language from file extension
                        ext = file_name.split(".")[-1] if "." in file_name else ""
                        st.code(file_content, language=ext)
                        if st.button("Close Preview", key=f"close_preview_{idx}"):
                            st.session_state[f"preview_{file_path}"] = False
                            st.rerun()
                else:
                    st.error("Failed to load file preview. File may be binary or too large.")
                    if st.button("Close", key=f"close_preview_err_{idx}"):
                        st.session_state[f"preview_{file_path}"] = False
                        st.rerun()
                st.markdown("---")

def import_github_file_to_chat_context(content: str, filename: str, repo_fullname: str):
    st.session_state.cb_file_text = content
    st.session_state.cb_img_bytes = None
    st.session_state.cb_img_mime = None
    st.session_state.cb_df = None
    st.session_state.cb_df_name = ""
    
    # Save attachment record
    save_attachment_record(filename, "text/plain", source="github")
    
    ensure_active_conversation()
    
    user_msg = f"Imported file: {filename} from repository {repo_fullname}"
    assistant_msg = f"📄 **{filename}** has been successfully imported from GitHub ({repo_fullname}) and attached to this conversation. You can now ask questions about the code!"
    
    st.session_state.cb_messages.append({
        "role": "user",
        "content": user_msg,
        "snippet": filename
    })
    st.session_state.cb_messages.append({
        "role": "assistant",
        "content": assistant_msg
    })
    
    if db_enabled() and st.session_state.active_conversation_id:
        db_action(save_message, int(st.session_state.active_conversation_id), "user", user_msg, filename)
        db_action(save_message, int(st.session_state.active_conversation_id), "assistant", assistant_msg)

def start_new_chat():
    sync_active_conversation()
    st.session_state.cb_messages = []
    st.session_state.cb_file_text = ""
    st.session_state.cb_img_bytes = None
    st.session_state.cb_img_mime = None
    st.session_state.cb_df = None
    st.session_state.cb_df_name = ""
    st.session_state.ai_insights_narrative = ""
    st.session_state.up_keys["doc"] += 1
    st.session_state.up_keys["img"] += 1
    st.session_state.up_keys["data"] += 1
    st.session_state.active_conversation_id = None
    st.session_state.selected_nav = "Chat"

    st.session_state.last_synced_id = None
    st.session_state.last_synced_title = ""
    st.session_state.last_synced_provider = ""
    st.session_state.last_synced_model = ""

def get_conversation_title(messages):
    for msg in messages:
        if msg.get("role") == "user":
            title = " ".join(str(msg.get("content", "")).split())
            if title:
                return title[:29] + "..." if len(title) > 32 else title
    return "New chat"

def ensure_active_conversation():
    if st.session_state.active_conversation_id:
        return st.session_state.active_conversation_id

    if db_enabled():
        conversation_id = create_conversation(
            st.session_state.db_user_id,
            provider=st.session_state.get("llm_provider", ""),
            model_name=st.session_state.get("active_model_name", ""),
        )
        st.session_state.active_conversation_id = conversation_id
        return conversation_id

    st.session_state.conversation_counter += 1
    conversation_id = f"chat_{st.session_state.conversation_counter}"
    st.session_state.active_conversation_id = conversation_id
    timestamp = datetime.datetime.now().isoformat()
    st.session_state.conversations.append({
        "id": conversation_id,
        "title": get_conversation_title(st.session_state.cb_messages),
        "messages": list(st.session_state.cb_messages),
        "cb_df": st.session_state.cb_df,
        "cb_df_name": st.session_state.cb_df_name,
        "cb_file_text": st.session_state.cb_file_text,
        "cb_img_bytes": st.session_state.cb_img_bytes,
        "cb_img_mime": st.session_state.cb_img_mime,
        "created_at": timestamp,
        "updated_at": timestamp,
    })
    return conversation_id

def sync_active_conversation():
    if not st.session_state.cb_messages and not st.session_state.active_conversation_id:
        return

    conversation_id = ensure_active_conversation()
    title = get_conversation_title(st.session_state.cb_messages)
    provider = st.session_state.get("llm_provider", "")
    model_name = st.session_state.get("active_model_name", "")

    # Skip DB update if active conversation data hasn't changed since last sync
    if (
        st.session_state.get("last_synced_id") == conversation_id
        and st.session_state.get("last_synced_title") == title
        and st.session_state.get("last_synced_provider") == provider
        and st.session_state.get("last_synced_model") == model_name
    ):
        return

    if db_enabled():
        db_action(
            update_conversation,
            conversation_id,
            title,
            provider,
            model_name,
        )
        st.session_state.last_synced_id = conversation_id
        st.session_state.last_synced_title = title
        st.session_state.last_synced_provider = provider
        st.session_state.last_synced_model = model_name
        return

    for conversation in st.session_state.conversations:
        if conversation["id"] == conversation_id:
            conversation["messages"] = list(st.session_state.cb_messages)
            conversation["title"] = title
            conversation["cb_df"] = st.session_state.cb_df
            conversation["cb_df_name"] = st.session_state.cb_df_name
            conversation["cb_file_text"] = st.session_state.cb_file_text
            conversation["cb_img_bytes"] = st.session_state.cb_img_bytes
            conversation["cb_img_mime"] = st.session_state.cb_img_mime
            conversation["ai_insights_narrative"] = st.session_state.get("ai_insights_narrative", "")
            conversation["updated_at"] = datetime.datetime.now().isoformat()
            break

    st.session_state.last_synced_id = conversation_id
    st.session_state.last_synced_title = title
    st.session_state.last_synced_provider = provider
    st.session_state.last_synced_model = model_name

def load_conversation(conversation_id):
    sync_active_conversation()
    
    # Reset current attachments first
    st.session_state.cb_df = None
    st.session_state.cb_df_name = ""
    st.session_state.cb_file_text = ""
    st.session_state.cb_img_bytes = None
    st.session_state.cb_img_mime = None
    st.session_state.ai_insights_narrative = ""

    if db_enabled():
        st.session_state.active_conversation_id = int(conversation_id)
        st.session_state.cb_messages = cached_load_messages(int(conversation_id)) or []
        st.session_state.selected_nav = "Chat"

        # Load and restore associated attachments from database
        uploads = cached_load_conversation_uploads(int(conversation_id)) or []
        has_dataset = False
        for up in uploads:
            ftype = up["file_type"]
            if ftype == "dataset" and up["data_json"]:
                try:
                    import pandas as pd
                    st.session_state.cb_df = pd.read_json(up["data_json"])
                    st.session_state.cb_df_name = up["filename"]
                    has_dataset = True
                except Exception:
                    pass
            elif ftype == "image" and up["image_data"]:
                st.session_state.cb_img_bytes = up["image_data"]
                st.session_state.cb_img_mime = up["mime_type"]
            elif ftype == "document" and up["text_content"]:
                st.session_state.cb_file_text = up["text_content"]

        if has_dataset:
            auto_generate_insights()

        st.session_state.last_synced_id = int(conversation_id)
        st.session_state.last_synced_title = get_conversation_title(st.session_state.cb_messages)
        st.session_state.last_synced_provider = st.session_state.get("llm_provider", "")
        st.session_state.last_synced_model = st.session_state.get("active_model_name", "")
        return

    for conversation in st.session_state.conversations:
        if conversation["id"] == conversation_id:
            st.session_state.active_conversation_id = conversation_id
            st.session_state.cb_messages = list(conversation.get("messages", []))
            st.session_state.selected_nav = "Chat"

            # Restore attachments from local memory
            st.session_state.cb_df = conversation.get("cb_df")
            st.session_state.cb_df_name = conversation.get("cb_df_name", "")
            st.session_state.cb_file_text = conversation.get("cb_file_text", "")
            st.session_state.cb_img_bytes = conversation.get("cb_img_bytes")
            st.session_state.cb_img_mime = conversation.get("cb_img_mime")
            st.session_state.ai_insights_narrative = conversation.get("ai_insights_narrative", "")

            # If dataset is loaded but narrative is missing, generate it
            if st.session_state.cb_df is not None and not st.session_state.ai_insights_narrative:
                auto_generate_insights()

            st.session_state.last_synced_id = conversation_id
            st.session_state.last_synced_title = get_conversation_title(st.session_state.cb_messages)
            st.session_state.last_synced_provider = st.session_state.get("llm_provider", "")
            st.session_state.last_synced_model = st.session_state.get("active_model_name", "")
            return

def get_conversation_order(conversation):
    conversation_id = str(conversation.get("id", ""))
    if conversation_id.startswith("chat_"):
        try:
            return int(conversation_id.split("_", 1)[1])
        except ValueError:
            pass
    return 0

def get_recent_conversations(limit=1000):
    sync_active_conversation()
    if db_enabled():
        if st.session_state.get("recent_conversations") is None:
            st.session_state.recent_conversations = cached_list_conversations(st.session_state.db_user_id, 1000) or []
        return st.session_state.recent_conversations[:limit]

    conversations = [
        conversation for conversation in st.session_state.conversations
        if conversation.get("messages")
    ]
    return sorted(conversations, key=get_conversation_order, reverse=True)[:limit]

def is_vision_model(provider, model):
    if not model:
        return False
    model_l = model.lower()
    if provider == "Gemini":
        return True  # All Gemini models generally support vision or we assume they do
    if provider == "Groq":
        return "vision" in model_l or "llama-3.2" in model_l
    if provider == "BazaarLink":
        return "vision" in model_l or "gpt-4" in model_l or "claude" in model_l or "gemini" in model_l
    return False

def build_openai_messages(sys_prompt, previous_messages, current_text, image_bytes=None, image_mime=None, is_vision=True):
    messages = [{"role": "system", "content": sys_prompt}]
    for msg in previous_messages:
        role = "user" if msg["role"] == "user" else "assistant"
        messages.append({"role": role, "content": msg["content"]})

    if image_bytes and is_vision:
        mime = image_mime or "image/jpeg"
        encoded = base64.b64encode(image_bytes).decode("utf-8")
        messages.append({
            "role": "user",
            "content": [
                {"type": "text", "text": current_text},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{encoded}"}},
            ],
        })
    else:
        messages.append({"role": "user", "content": current_text})

    return messages

MODEL_OPTIONS = {
    "Gemini": {
        "title": "Gemini",
        "subtitle": "Auto-selects an available Gemini model for your key",
        "provider": "Gemini",
        "pill": "Gemini",
    },
    "Groq": {
        "title": "Groq",
        "subtitle": "Auto-selects an available Groq model for your key",
        "provider": "Groq",
        "pill": "Groq",
    },
    "BazaarLink": {
        "title": "BazaarLink",
        "subtitle": "Unified access to AI models (OpenAI compatible)",
        "provider": "BazaarLink",
        "pill": "BazaarLink",
    },
}

GEMINI_MODEL_PRIORITY = [
    "gemini-2.5-flash",
    "gemini-2.0-flash",
    "gemini-1.5-flash",
    "gemini-1.5-pro",
]

@st.cache_data(show_spinner=False, ttl=3600)
def get_available_groq_models(api_key):
    if not api_key or api_key == "YOUR_GROQ_API_KEY_HERE":
        return [], "Groq API key is missing or is set to placeholder."

    OpenAI = import_openai_client()
    if OpenAI is None:
        return [], "OpenAI package is not installed."

    try:
        client = OpenAI(base_url="https://api.groq.com/openai/v1", api_key=api_key, timeout=10.0)
        model_ids = [model.id for model in client.models.list().data]
        return sorted(set(model_ids)), ""
    except Exception as e:
        fallback = [
            "llama-3.3-70b-versatile",
            "llama-3.1-8b-instant",
            "mixtral-8x7b-32768",
            "gemma2-9b-it"
        ]
        return fallback, str(e)

@st.cache_data(show_spinner=False, ttl=3600)
def get_available_gemini_models(api_key):
    if not api_key or api_key == "YOUR_GEMINI_API_KEY_HERE":
        return [], "Gemini API key is missing or is set to placeholder."

    genai = import_genai_module()
    if genai is None:
        return [], "google-generativeai package is not installed."

    try:
        genai.configure(api_key=api_key)
        available = []
        for model in genai.list_models():
            methods = getattr(model, "supported_generation_methods", [])
            if "generateContent" in methods:
                name = model.name.split("/")[-1]
                if name.startswith("gemini"):
                    available.append(name)
        return sorted(set(available)), ""
    except Exception as e:
        fallback = [
            "gemini-2.5-flash",
            "gemini-2.0-flash",
            "gemini-1.5-flash",
            "gemini-1.5-pro",
        ]
        return fallback, str(e)

@st.cache_data(show_spinner=False, ttl=3600)
def get_available_bazaarlink_models(api_key):
    if not api_key:
        return [], "BazaarLink API key is missing. Check BAZAARLINK_API_KEY environment variable."

    OpenAI = import_openai_client()
    if OpenAI is None:
        return [], "OpenAI package is not installed."

    try:
        client = OpenAI(
            base_url=BAZAARLINK_BASE_URL,
            api_key=api_key,
            default_headers=BAZAARLINK_HEADERS,
            timeout=10.0
        )
        model_ids = [model.id for model in client.models.list().data]
        return sorted(set(model_ids)), ""
    except Exception as e:
        fallback = [
            "auto:free",
            "gpt-4o",
            "gpt-4o-mini",
            "claude-3-5-sonnet",
            "gemini-1.5-flash",
            "deepseek-chat"
        ]
        err_msg = str(e)
        if "cloudflare" in err_msg.lower() or "<html" in err_msg.lower() or "doctype" in err_msg.lower() or "just a moment" in err_msg.lower():
            err_msg = "Cloudflare security challenge or block encountered"
        elif len(err_msg) > 100:
            err_msg = err_msg[:97] + "..."
        return fallback, err_msg

def choose_available_model(available_models, preferred_models):
    if not available_models:
        return ""

    for preferred in preferred_models:
        if preferred in available_models:
            return preferred

    return available_models[0]

def get_active_model(provider):
    try:
        if provider == "Gemini":
            available, err = get_available_gemini_models(GEMINI_API_KEY)
            return choose_available_model(available, GEMINI_MODEL_PRIORITY), available, err

        if provider == "Groq":
            available, err = get_available_groq_models(GROQ_API_KEY)
            return choose_available_model(available, GROQ_MODEL_PRIORITY), available, err

        if provider == "BazaarLink":
            available, err = get_available_bazaarlink_models(BAZAARLINK_API_KEY)
            return choose_available_model(available, BAZAARLINK_MODEL_PRIORITY), available, err

        raise ValueError(f"Unknown or unsupported provider: {provider}")
    except Exception as ex:
        return "", [], str(ex)

def get_voice_transcription_model():
    active_model, _, model_error = get_active_model("Gemini")
    if active_model:
        return active_model, ""
    return "", model_error or "No compatible Gemini model was found for voice transcription."

def render_model_picker():
    selected_key = st.session_state.llm_provider
    if selected_key not in MODEL_OPTIONS:
        selected_key = "Gemini"
        st.session_state.llm_provider = selected_key

    selected = MODEL_OPTIONS[selected_key]
    active_model, available_models, model_error = get_active_model(selected["provider"])
    
    # Check if a model is already selected in session state for the active provider
    if selected["provider"] == "Gemini":
        if not st.session_state.get("cb_model") or st.session_state.cb_model not in available_models:
            st.session_state.cb_model = active_model
        st.session_state.active_model_name = st.session_state.cb_model
    elif selected["provider"] == "Groq":
        if not st.session_state.get("groq_model") or st.session_state.groq_model not in available_models:
            st.session_state.groq_model = active_model
        st.session_state.active_model_name = st.session_state.groq_model
    elif selected["provider"] == "BazaarLink":
        if not st.session_state.get("bazaarlink_model") or st.session_state.bazaarlink_model not in available_models:
            st.session_state.bazaarlink_model = active_model
        st.session_state.active_model_name = st.session_state.bazaarlink_model

    _, search_toggle_col, picker_col = st.columns([3.8, 1.6, 1.6])
    with search_toggle_col:
        st.toggle("🌐 Web Search", key="web_search_enabled", help="Enable DuckDuckGo search for real-time information")
    with picker_col:
        pill_label = f"{selected['pill']} ({st.session_state.active_model_name})" if st.session_state.active_model_name else selected["pill"]
        with st.popover(pill_label, use_container_width=True):
            if model_error:
                st.warning(f"⚠️ API Info: {model_error}. Using fallback list.")
            st.markdown('<div class="model-menu-title">Select Platform</div>', unsafe_allow_html=True)
            for key, option in MODEL_OPTIONS.items():
                active = "✓ " if key == selected_key else ""
                label = option["title"]
                if st.button(label, key=f"llm_choice_{key}", use_container_width=True):
                    st.session_state.llm_provider = key
                    st.session_state.active_model_name = ""  # reset to default on provider change
                    if key == "Gemini":
                        st.session_state.cb_model = ""
                    elif key == "Groq":
                        st.session_state.groq_model = ""
                    elif key == "BazaarLink":
                        st.session_state.bazaarlink_model = ""
                    st.rerun()
            
            if available_models:
                st.markdown('<div class="model-menu-title" style="margin-top:12px;">Select Model</div>', unsafe_allow_html=True)
                
                # Determine default index
                try:
                    default_idx = available_models.index(st.session_state.active_model_name)
                except ValueError:
                    default_idx = 0
                    
                selected_model = st.selectbox(
                    "Model Options",
                    available_models,
                    index=default_idx,
                    key="active_model_selection_selectbox",
                    label_visibility="collapsed"
                )
                
                if selected_model != st.session_state.active_model_name:
                    st.session_state.active_model_name = selected_model
                    if selected["provider"] == "Gemini":
                        st.session_state.cb_model = selected_model
                    elif selected["provider"] == "Groq":
                        st.session_state.groq_model = selected_model
                    elif selected["provider"] == "BazaarLink":
                        st.session_state.bazaarlink_model = selected_model
                    st.rerun()

def render_sidebar():
    selected = MODEL_OPTIONS.get(st.session_state.llm_provider, MODEL_OPTIONS["Gemini"])
    active_model = st.session_state.active_model_name or "Checking available models"
    with st.sidebar:
        st.markdown("""
        <div class="side-brand">
            <div class="side-logo">
                <i class="fa-solid fa-robot" style="font-size: 1.45rem; color: #da7756;"></i>
            </div>
            <div>
                <h3>Morepen AI</h3>
                <p>Assistant workspace</p>
            </div>
        </div>
        """, unsafe_allow_html=True)
        col1, col2 = st.columns(2)
        with col1:
            if st.button("New Chat", key="nav_Chat", use_container_width=True):
                start_new_chat()
                st.rerun()
        with col2:
            with st.popover("Delete", use_container_width=True):
                st.markdown('<div style="color:#ef4444; font-weight:700; margin-bottom:8px; font-size:0.9rem;">Delete Chat</div>', unsafe_allow_html=True)
                recent_chats = get_recent_conversations(limit=30)
                if recent_chats:
                    chat_list = []
                    chat_options = {}
                    for idx, c in enumerate(recent_chats):
                        label = f"{idx+1}. {c['title']}"
                        chat_list.append(label)
                        chat_options[label] = c["id"]
                    
                    selected_chat_label = st.selectbox("Choose chat", chat_list, key="delete_chat_select")
                    selected_chat_id = chat_options[selected_chat_label]
                    if st.button("Confirm", type="primary", key="confirm_delete_chat", use_container_width=True):
                        if db_enabled():
                            db_action(delete_conversation, selected_chat_id)
                        else:
                            st.session_state.conversations = [c for c in st.session_state.conversations if c["id"] != selected_chat_id]
                        
                        # If deleting the currently active conversation, start a new chat
                        if st.session_state.active_conversation_id == selected_chat_id:
                            start_new_chat()
                        st.success("Deleted!")
                        st.rerun()
                else:
                    st.info("No chats to delete.")

        st.markdown('<div class="side-section-title">Recents</div>', unsafe_allow_html=True)
        # Fetch only up to current limit for displaying
        recent_conversations = get_recent_conversations(limit=st.session_state.chats_limit)
        
        search_query = st.text_input("Search chats", key="chat_search", label_visibility="collapsed", placeholder="Search chats...")
        if search_query:
            # For search query, search all conversations (up to 1000)
            all_recent = get_recent_conversations(limit=1000)
            recent_conversations = [c for c in all_recent if search_query.lower() in c["title"].lower()]
            
        if recent_conversations:
            recents_container = st.container(height=260, border=False)
            with recents_container:
                for conversation in recent_conversations:
                    if st.button(conversation["title"], key=f'recent_{conversation["id"]}', use_container_width=True):
                        load_conversation(conversation["id"])
                        st.rerun()
                
                # Render "Load More" button if there are more chats to display
                total_available = 0
                if db_enabled():
                    all_cached = st.session_state.get("recent_conversations")
                    if all_cached:
                        total_available = len(all_cached)
                else:
                    total_available = len([c for c in st.session_state.conversations if c.get("messages")])
                
                if total_available > st.session_state.chats_limit and not search_query:
                    if st.button("Load More ➕", key="load_more_chats", use_container_width=True):
                        st.session_state.chats_limit += 30
                        st.rerun()
        else:
            if search_query:
                st.info("No matching chats.")
            else:
                if st.button("No recent chats yet", key="recent_empty", use_container_width=True, disabled=True):
                    st.rerun()

        st.markdown('<div class="side-section-title">Additional Features</div>', unsafe_allow_html=True)
        if st.button("Upload from Google Drive", key="feature_Google_Drive_Upload", use_container_width=True):
            st.session_state.selected_nav = "Google Drive Upload"
            st.rerun()
            
        if st.button("Data Analysis Workspace", key="feature_Data_Analysis", use_container_width=True):
            st.session_state.selected_nav = "Data Analysis"
            st.rerun()

        if st.button("Document & File Library", key="feature_Document_Library", use_container_width=True):
            st.session_state.selected_nav = "Document Library"
            st.rerun()


        # Load GitHub Credentials to see if Explorer should be shown
        github_creds = st.session_state.get("github_credentials")
        if github_creds is None:
            if db_enabled() and st.session_state.db_user_id:
                from database import load_github_credentials
                db_creds = db_action(load_github_credentials, st.session_state.db_user_id)
                if db_creds:
                    github_creds = db_creds
                    st.session_state.github_credentials = db_creds
                else:
                    st.session_state.github_credentials = False
                    github_creds = False
            else:
                st.session_state.github_credentials = False
                github_creds = False
        
        github_linked = bool(github_creds)
        if github_linked:
            if st.button("GitHub Repository Explorer", key="feature_Github_Explorer", use_container_width=True):
                st.session_state.selected_nav = "GitHub Explorer"
                st.rerun()

        # --- Google integration status and controls ---
        st.markdown('<div class="side-section-title">Google Integration</div>', unsafe_allow_html=True)
        
        # Load Google Credentials (cached in session state)
        google_creds = st.session_state.get("google_credentials")
        if google_creds is None:
            if db_enabled() and st.session_state.db_user_id:
                from database import load_google_credentials
                db_creds = db_action(load_google_credentials, st.session_state.db_user_id)
                if db_creds:
                    google_creds = db_creds
                    st.session_state.google_credentials = db_creds
                else:
                    st.session_state.google_credentials = False
                    google_creds = False
            else:
                st.session_state.google_credentials = False
                google_creds = False

        google_linked = bool(google_creds)
        if google_linked:
            st.markdown(
                """
                <div style="background: rgba(16, 185, 129, 0.1); border: 1px solid rgba(16, 185, 129, 0.3); border-radius: 10px; padding: 10px; margin-bottom: 10px;">
                    <span style="color: #10b981; font-weight: 700; font-size: 0.85rem;">🟢 Google Services Connected</span>
                </div>
                """,
                unsafe_allow_html=True
            )
            if st.button("Disconnect Google Account", key="disconnect_google_btn", use_container_width=True):
                st.session_state.google_credentials = False
                if db_enabled() and st.session_state.db_user_id:
                    from database import delete_google_credentials
                    db_action(delete_google_credentials, st.session_state.db_user_id)
                st.success("Disconnected Google Account!")
                st.rerun()
        else:
            google_oauth_scopes = (
                "https://www.googleapis.com/auth/gmail.compose "
                "https://www.googleapis.com/auth/calendar.events "
                "https://www.googleapis.com/auth/documents "
                "https://www.googleapis.com/auth/drive.file "
                "openid email profile"
            )
            import urllib.parse
            google_auth_url = (
                f"https://accounts.google.com/o/oauth2/v2/auth?"
                f"client_id={GOOGLE_CLIENT_ID}&"
                f"redirect_uri={urllib.parse.quote(GOOGLE_REDIRECT_URI)}&"
                f"response_type=code&"
                f"scope={urllib.parse.quote(google_oauth_scopes)}&"
                f"access_type=offline&"
                f"prompt=consent&"
                f"state=connect_google"
            )
            st.markdown(f'<a href="{google_auth_url}" style="text-decoration:none; cursor: pointer !important; pointer-events: auto !important;"><button style="width:100%; height:38px; margin-bottom:10px; border-radius:10px; border:1px solid #da7756; background:#da7756; color:white; font-weight:700; cursor:pointer;" onmouseover="this.style.background=\'#c56241\'" onmouseout="this.style.background=\'#da7756\'">🔗 Connect Google Services</button></a>', unsafe_allow_html=True)

        # --- GitHub integration status and controls ---
        st.markdown('<div class="side-section-title">GitHub Integration</div>', unsafe_allow_html=True)
        if github_linked:
            st.markdown(
                """
                <div style="background: rgba(16, 185, 129, 0.1); border: 1px solid rgba(16, 185, 129, 0.3); border-radius: 10px; padding: 10px; margin-bottom: 10px;">
                    <span style="color: #10b981; font-weight: 700; font-size: 0.85rem;">🟢 GitHub Connected</span>
                </div>
                """,
                unsafe_allow_html=True
            )
            if st.button("Disconnect GitHub Account", key="disconnect_github_btn", use_container_width=True):
                st.session_state.github_credentials = False
                if db_enabled() and st.session_state.db_user_id:
                    from database import delete_github_credentials
                    db_action(delete_github_credentials, st.session_state.db_user_id)
                st.success("Disconnected GitHub Account!")
                st.rerun()
        else:
            if GITHUB_CLIENT_ID and GITHUB_CLIENT_SECRET:
                import urllib.parse
                github_oauth_scope = "repo read:user"
                signed_state = generate_signed_github_state(st.session_state.get("username", ""))
                github_auth_url = (
                    f"https://github.com/login/oauth/authorize?"
                    f"client_id={GITHUB_CLIENT_ID}&"
                    f"redirect_uri={urllib.parse.quote(GOOGLE_REDIRECT_URI)}&"
                    f"scope={urllib.parse.quote(github_oauth_scope)}&"
                    f"state={urllib.parse.quote(signed_state)}"
                )
                st.markdown(f'<a href="{github_auth_url}" style="text-decoration:none; cursor: pointer !important; pointer-events: auto !important;"><button style="width:100%; height:38px; margin-bottom:10px; border-radius:10px; border:1px solid #24292e; background:#24292e; color:white; font-weight:700; cursor:pointer;" onmouseover="this.style.background=\'#444d56\'" onmouseout="this.style.background=\'#24292e\'">🔗 Connect GitHub</button></a>', unsafe_allow_html=True)
            else:
                st.caption("ℹ️ Configure GITHUB_CLIENT_ID and GITHUB_CLIENT_SECRET to enable GitHub connection.")



        st.markdown('<div class="side-section-title">Current model</div>', unsafe_allow_html=True)
        st.markdown(
            f"""
            <div class="side-status" style="margin-top:0;">
                <strong>{selected["title"]}</strong><br>
                <span>{active_model}</span>
            </div>
            """,
            unsafe_allow_html=True
        )

        st.markdown('<div class="side-section-title">Session</div>', unsafe_allow_html=True)
        current_user = st.session_state.get("username", "Guest")
        st.markdown(
            f"""
            <div class="side-status" style="margin-top:0; margin-bottom:10px; padding:10px 12px; display:flex; align-items:center; gap:10px;">
                <i class="fa-solid fa-circle-user" style="font-size: 1.4rem; color: #da7756;"></i>
                <div>
                    <span style="font-size:0.75rem; color:#6b685c; text-transform:uppercase; font-weight:800; letter-spacing:0.05em; display:block; line-height:1;">Signed in as</span>
                    <strong style="color:#1c1b1a; font-size:0.88rem; display:block; margin-top:3px; word-break:break-all;">{current_user}</strong>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )
        if st.button("Logout", use_container_width=True):
            st.session_state.authenticated = False
            st.session_state.username = ""
            st.session_state.db_user_id = None
            st.session_state.cb_messages = []
            st.session_state.active_conversation_id = None
            st.session_state.cb_df = None
            st.session_state.cb_df_name = ""
            st.session_state.cb_file_text = ""
            st.session_state.cb_img_bytes = None
            st.session_state.cb_img_mime = None
            st.session_state.ai_insights_narrative = ""
            st.session_state.recent_conversations = None
            st.rerun()



# --- detect model once ---
if st.session_state.cb_model is None or st.session_state.cb_model == "gpt-4o":
    st.session_state.cb_model = "gemini-1.5-flash"

record_profiler_checkpoint("Pre-Sidebar Setup")
render_sidebar()
record_profiler_checkpoint("Sidebar Rendering")

if st.session_state.selected_nav == "Search Chats":
    st.subheader("Search Chats")
    
    @safe_fragment
    def search_chats_block():
        query = st.text_input("Search chat history", placeholder="Type a keyword", key="search_chats_keyword")
        if query:
            matches = [
                msg for msg in st.session_state.cb_messages
                if query.lower() in str(msg.get("content", "")).lower()
            ]
            if matches:
                for msg in matches:
                    label = "You" if msg["role"] == "user" else "Assistant"
                    with st.expander(label):
                        st.markdown(msg["content"])
            else:
                st.info("No matching chats found.")
        else:
            st.info("Type a keyword to search your current chat history.")
            
    search_chats_block()
    record_profiler_checkpoint("Search Chats View")
    st.stop()

if st.session_state.selected_nav == "GitHub Explorer":
    render_github_explorer()
    record_profiler_checkpoint("GitHub Explorer View")
    st.stop()

if st.session_state.selected_nav == "Documents":
    st.subheader("Documents")
    if st.session_state.cb_file_text:
        st.text_area("Current document text", st.session_state.cb_file_text, height=420)
    else:
        st.info("No document is attached yet. Go to Chat and upload a PDF, TXT, DOC, DOCX, or PPTX file with your message.")
    record_profiler_checkpoint("Documents View")
    st.stop()

if st.session_state.selected_nav == "Data Analysis":
    render_data_analysis()
    record_profiler_checkpoint("Data Analysis View")
    st.stop()

if st.session_state.selected_nav == "Google Drive Upload":
    render_google_drive_upload()
    record_profiler_checkpoint("Google Drive Upload View")
    st.stop()

if st.session_state.selected_nav == "Document Library":
    render_document_library()
    record_profiler_checkpoint("Document Library View")
    st.stop()

if st.session_state.selected_nav == "History":
    st.subheader("History")
    if st.session_state.cb_messages:
        for msg in st.session_state.cb_messages:
            label = "You" if msg["role"] == "user" else "Assistant"
            with st.expander(label):
                st.markdown(msg["content"])
    else:
        st.info("No chat history yet.")
    record_profiler_checkpoint("History View")
    st.stop()

# --- toolbar ---
tc1, tc2, tc3 = st.columns([6, 1, 1])
with tc2:
    if st.button("🗑️ Clear", use_container_width=True):
        if db_enabled() and st.session_state.active_conversation_id:
            db_action(clear_conversation_messages, int(st.session_state.active_conversation_id))
        st.session_state.cb_messages = []
        st.session_state.cb_file_text = ""
        st.session_state.cb_img_bytes = None
        st.session_state.cb_img_mime = None
        st.session_state.cb_df = None
        st.session_state.cb_df_name = ""
        st.session_state.ai_insights_narrative = ""
        st.session_state.up_keys["doc"] += 1
        st.session_state.up_keys["img"] += 1
        st.session_state.up_keys["data"] += 1
        st.rerun()
def generate_docx_transcript(messages):
    import docx
    import io
    
    doc = docx.Document()
    doc.add_heading('Morepen AI - Chat Transcript', level=1)
    ist_tz = datetime.timezone(datetime.timedelta(hours=5, minutes=30))
    now_ist = datetime.datetime.now(ist_tz)
    doc.add_paragraph('Exported on: ' + now_ist.strftime('%Y-%m-%d %H:%M:%S'))
    doc.add_paragraph()
    
    for msg in messages:
        role_label = "You" if msg["role"] == "user" else "Assistant"
        p_role = doc.add_paragraph()
        r_run = p_role.add_run(f"{role_label}:")
        r_run.bold = True
        
        doc.add_paragraph(msg["content"])
        doc.add_paragraph()
        
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def generate_pdf_transcript(messages):
    import io
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib import colors
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    story = []
    
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=20,
        leading=24,
        textColor=colors.HexColor('#4f46e5'),
        spaceAfter=10
    )
    
    meta_style = ParagraphStyle(
        'DocMeta',
        parent=styles['Normal'],
        fontSize=9,
        textColor=colors.HexColor('#6b7280'),
        spaceAfter=15
    )
    
    user_label_style = ParagraphStyle(
        'UserLabel',
        parent=styles['Heading2'],
        fontSize=11,
        leading=14,
        textColor=colors.HexColor('#06b6d4'),
        spaceBefore=8,
        spaceAfter=2
    )
    
    assistant_label_style = ParagraphStyle(
        'AssistantLabel',
        parent=styles['Heading2'],
        fontSize=11,
        leading=14,
        textColor=colors.HexColor('#da7756'),
        spaceBefore=8,
        spaceAfter=2
    )
    
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#1f2937'),
        spaceAfter=10
    )
    
    ist_tz = datetime.timezone(datetime.timedelta(hours=5, minutes=30))
    now_ist = datetime.datetime.now(ist_tz)
    story.append(Paragraph("Morepen AI - Chat Transcript", title_style))
    story.append(Paragraph(f"Exported on: {now_ist.strftime('%Y-%m-%d %H:%M:%S')}", meta_style))
    story.append(Spacer(1, 10))
    
    for msg in messages:
        role = msg["role"]
        label = "You" if role == "user" else "Assistant"
        label_style = user_label_style if role == "user" else assistant_label_style
        
        story.append(Paragraph(label, label_style))
        content = msg["content"]
        content_escaped = content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")
        story.append(Paragraph(content_escaped, body_style))
        story.append(Spacer(1, 4))
        
    doc.build(story)
    return buffer.getvalue()

with tc3:
    with st.popover("💾 Export", use_container_width=True):
        st.markdown("<div style='font-size:0.9rem; font-weight:700; margin-bottom:8px;'>Export Format</div>", unsafe_allow_html=True)
        
        ist_tz = datetime.timezone(datetime.timedelta(hours=5, minutes=30))
        now_ist = datetime.datetime.now(ist_tz)
        
        # 1. Plain Text Download
        lines = []
        for msg in st.session_state.cb_messages:
            label = "You" if msg["role"] == "user" else "Assistant"
            lines.append(f"[{label}]\n{msg['content']}\n")
        export_txt = "\n".join(lines) if lines else "No messages to export yet."
        
        st.download_button(
            "TXT Format (.txt)",
            data=export_txt,
            file_name=f"chat_{now_ist.strftime('%Y%m%d_%H%M%S')}.txt",
            mime="text/plain",
            use_container_width=True
        )
        
        # Disable Word and PDF downloads if no messages
        has_messages = len(st.session_state.cb_messages) > 0
        
        # 2. Word Document Download
        if has_messages:
            try:
                docx_data = generate_docx_transcript(st.session_state.cb_messages)
                st.download_button(
                    "Word Format (.docx)",
                    data=docx_data,
                    file_name=f"chat_{now_ist.strftime('%Y%m%d_%H%M%S')}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            except Exception as e:
                st.error(f"Docx export error: {e}")
        else:
            st.button("Word Format (.docx)", disabled=True, use_container_width=True)
            
        # 3. PDF Document Download
        if has_messages:
            try:
                pdf_data = generate_pdf_transcript(st.session_state.cb_messages)
                st.download_button(
                    "PDF Format (.pdf)",
                    data=pdf_data,
                    file_name=f"chat_{now_ist.strftime('%Y%m%d_%H%M%S')}.pdf",
                    mime="application/pdf",
                    use_container_width=True
                )
            except Exception as e:
                st.error(f"PDF export error: {e}")
        else:
            st.button("PDF Format (.pdf)", disabled=True, use_container_width=True)

        # --- Google Docs Export divider ---
        st.markdown(
            "<hr style='margin:12px 0 10px 0; border:none; border-top:1px solid #e5e3d9;'>"
            "<div style='font-size:0.9rem; font-weight:700; margin-bottom:8px;'>Google Docs Export</div>",
            unsafe_allow_html=True
        )

        # 5. Google Docs Export
        google_creds_export = st.session_state.get("google_credentials")
        gdocs_access_token = None
        if google_creds_export:
            try:
                from backend.google_service import get_valid_token
                gdocs_access_token, _updated = get_valid_token(
                    st.session_state.get("db_user_id"), google_creds_export
                )
                if _updated:
                    st.session_state.google_credentials = _updated
            except Exception:
                gdocs_access_token = None

        gdocs_connected = bool(gdocs_access_token)

        if gdocs_connected and has_messages:
            if st.button("Export to Google Docs", use_container_width=True, key="btn_export_gdocs"):
                from backend.export_service import export_to_google_docs, check_google_docs_scope
                has_docs_scope = check_google_docs_scope(gdocs_access_token)
                if not has_docs_scope:
                    # Prompt re-authorization with documents scope
                    if GOOGLE_CLIENT_ID:
                        import urllib.parse
                        docs_scope = (
                            "openid email profile "
                            "https://www.googleapis.com/auth/gmail.compose "
                            "https://www.googleapis.com/auth/calendar "
                            "https://www.googleapis.com/auth/documents "
                            "https://www.googleapis.com/auth/drive.file"
                        )
                        reauth_url = (
                            f"https://accounts.google.com/o/oauth2/v2/auth?"
                            f"client_id={GOOGLE_CLIENT_ID}&"
                            f"redirect_uri={urllib.parse.quote(GOOGLE_REDIRECT_URI)}&"
                            f"response_type=code&"
                            f"scope={urllib.parse.quote(docs_scope)}&"
                            f"state=connect_google&"
                            f"access_type=offline&"
                            f"prompt=consent"
                        )
                        st.warning(
                            "Google Docs export needs an additional permission. "
                            f"[Click here to re-authorize Google]({reauth_url}) "
                            "(takes a few seconds), then try again."
                        )
                    else:
                        st.error("Configure GOOGLE_CLIENT_ID to enable Google Docs export.")
                else:
                    conv_title = st.session_state.get("last_synced_title") or f"Chat {now_ist.strftime('%Y-%m-%d %H:%M')}"
                    with st.spinner("Creating Google Doc..."):
                        try:
                            gdoc_url = export_to_google_docs(
                                st.session_state.cb_messages,
                                title=conv_title,
                                access_token=gdocs_access_token,
                            )
                            st.session_state.export_status = ("gdocs", "success", gdoc_url)
                            st.success("✅ Exported to Google Docs!")
                            st.markdown(f"[🔗 Open Google Doc]({gdoc_url})")
                        except Exception as ex:
                            st.error(f"Google Docs export failed: {ex}")
        elif not gdocs_connected:
            st.button(
                "Export to Google Docs",
                disabled=True,
                use_container_width=True,
                help="Connect your Google account in the sidebar to enable Google Docs export.",
            )
            st.caption("ℹ️ Connect Google account in the sidebar to enable.")
        else:
            st.button("Export to Google Docs", disabled=True, use_container_width=True)

# Show export status notifications below the toolbar (persists across reruns)
if st.session_state.get("export_status"):
    _exp = st.session_state.export_status
    if _exp and len(_exp) == 3:
        _dest, _status, _url = _exp
        if _dest == "gdocs" and _status == "success":
            st.success(f"✅ Chat exported to Google Docs! [Open document]({_url})")
    st.session_state.export_status = None

# active badges
if st.session_state.drive_import_notice:
    st.success(st.session_state.drive_import_notice)
    st.session_state.drive_import_notice = ""
if st.session_state.cb_file_text:
    st.markdown('<span class="badge">📄 Document attached</span>', unsafe_allow_html=True)
if st.session_state.cb_img_bytes:
    st.markdown('<span class="badge">🖼️ Image attached</span>', unsafe_allow_html=True)
if st.session_state.cb_df is not None:
    st.markdown(f'<span class="badge">📊 {st.session_state.cb_df_name} attached</span>', unsafe_allow_html=True)

# --- quick prompts (data-aware) ---
has_df = st.session_state.cb_df is not None
if has_df:
    quick = ["📊 Summarize this dataset", "🔍 Find key insights", "📈 Identify trends",
             "🧹 Check data quality", "📉 Show correlations", "📋 Describe columns"]
else:
    quick = ["💡 Explain document", "📊 Summarize data", "🔍 Find insights",
             "✍️ Write report", "📈 Analyze trends", "❓ Help guide"]
st.markdown("**Quick prompts:**")
qcols = st.columns(len(quick))
chosen = None
for i, q in enumerate(quick):
    with qcols[i]:
        if st.button(q, key=f"cbq_{i}", use_container_width=True):
            chosen = q.split(" ", 1)[1]

record_profiler_checkpoint("Quick Prompts & Navigation Routing")

def render_copy_button(text, key):
    import urllib.parse
    escaped_text = urllib.parse.quote(text)
    button_html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
        <style>
            body {{
                margin: 0;
                padding: 0;
                overflow: hidden;
                background-color: transparent;
                display: flex;
                justify-content: center;
                align-items: center;
                height: 33px;
            }}
            .copy-btn {{
                display: inline-flex;
                align-items: center;
                justify-content: center;
                background-color: #fbfaf7;
                color: #6b685c;
                border: 1px solid #e5e3d9;
                font-size: 1rem;
                border-radius: 8px;
                cursor: pointer;
                transition: all 0.2s ease;
                user-select: none;
                height: 33px;
                width: 33px;
                box-sizing: border-box;
                margin-top: 0px;
            }}
            .copy-btn:hover {{
                background-color: #f5f2eb;
                border-color: #da7756;
                color: #da7756;
            }}
            .copy-btn:active {{
                transform: scale(0.95);
            }}
            .copied {{
                background-color: rgba(16, 185, 129, 0.08) !important;
                border-color: #10b981 !important;
                color: #10b981 !important;
            }}
        </style>
    </head>
    <body>
        <button id="btn" class="copy-btn" onclick="doCopy()">
            <i class="fa-regular fa-copy"></i>
        </button>

        <script>
            function doCopy() {{
                const text = decodeURIComponent('{escaped_text}');
                
                // Fallback sequence:
                // 1. Try parent navigator.clipboard
                try {{
                    if (window.parent && window.parent.navigator && window.parent.navigator.clipboard) {{
                        window.parent.navigator.clipboard.writeText(text).then(function() {{
                            showSuccess();
                        }}).catch(function(err) {{
                            tryIframeClipboard(text);
                        }});
                        return;
                    }}
                }} catch (e) {{
                    // CORS or other errors
                }}
                
                tryIframeClipboard(text);
            }}

            function tryIframeClipboard(text) {{
                // 2. Try iframe navigator.clipboard
                if (navigator.clipboard) {{
                    navigator.clipboard.writeText(text).then(function() {{
                        showSuccess();
                    }}).catch(function(err) {{
                        tryParentExecCommand(text);
                    }});
                }} else {{
                    tryParentExecCommand(text);
                }}
            }}

            function tryParentExecCommand(text) {{
                // 3. Try parent document execCommand (very reliable for same-origin iframe)
                try {{
                    if (window.parent && window.parent.document) {{
                        const pDoc = window.parent.document;
                        const el = pDoc.createElement('textarea');
                        el.value = text;
                        el.style.position = 'fixed';
                        el.style.left = '-9999px';
                        pDoc.body.appendChild(el);
                        el.select();
                        const successful = pDoc.execCommand('copy');
                        pDoc.body.removeChild(el);
                        if (successful) {{
                            showSuccess();
                            return;
                        }}
                    }}
                }} catch (e) {{
                    // CORS or other errors
                }}
                
                tryIframeExecCommand(text);
            }}

            function tryIframeExecCommand(text) {{
                // 4. Try iframe document execCommand
                const el = document.createElement('textarea');
                el.value = text;
                el.style.position = 'fixed';
                el.style.left = '-9999px';
                document.body.appendChild(el);
                el.select();
                try {{
                    const successful = document.execCommand('copy');
                    if (successful) {{
                        showSuccess();
                    }} else {{
                        console.error('execCommand copy failed');
                    }}
                }} catch (err) {{
                    console.error('All copy fallbacks failed: ', err);
                }}
                document.body.removeChild(el);
            }}

            function showSuccess() {{
                const btn = document.getElementById('btn');
                btn.classList.add('copied');
                btn.innerHTML = '<i class="fa-solid fa-check"></i>';
                setTimeout(function() {{
                    btn.classList.remove('copied');
                    btn.innerHTML = '<i class="fa-regular fa-copy"></i>';
                }}, 2000);
            }}
        </script>
    </body>
    </html>
    """
    import urllib.parse
    st.iframe(f"data:text/html;charset=utf-8,{urllib.parse.quote(button_html)}", height=33)

# --- chat history ---
chat_box = st.container()
with chat_box:
    for idx, msg in enumerate(st.session_state.cb_messages):
        with st.chat_message(msg["role"]):
            # Check if this is a file attachment/upload indicator message
            content = msg["content"]
            is_file_indicator = False
            filename = ""
            
            if content.startswith("Uploaded file: "):
                is_file_indicator = True
                filename = content[len("Uploaded file: "):]
            elif content.startswith("Restored file: "):
                is_file_indicator = True
                filename = content[len("Restored file: "):]
            elif content.startswith("Imported from Google Drive: "):
                is_file_indicator = True
                filename = content[len("Imported from Google Drive: "):]

            if is_file_indicator:
                # Render beautiful file badge instead of expander and text
                file_kind = get_file_type(filename)
                if file_kind == "dataset":
                    icon_class = "fa-solid fa-table"
                    icon_color = "#da7756" # Terracotta
                elif file_kind == "image":
                    icon_class = "fa-solid fa-image"
                    icon_color = "#ec4899" # Pink
                else:
                    icon_class = "fa-solid fa-file-lines"
                    icon_color = "#10b981" # Green
                
                st.markdown(
                    f'<div style="display: flex; align-items: center; gap: 12px; padding: 10px 14px; background-color: #fbf9f6; border: 1px solid #e5e3d9; border-radius: 8px; width: fit-content; margin: 4px 0 12px 0; box-shadow: 0 1px 3px rgba(0,0,0,0.05);">'
                    f'<i class="{icon_class}" style="color: {icon_color}; font-size: 1.5rem;"></i>'
                    f'<span style="font-weight: 700; color: #1c1b1a; font-size: 1.05rem;">{filename}</span>'
                    f'</div>',
                    unsafe_allow_html=True
                )
            else:
                if msg.get("img"):
                    try:
                        from PIL import Image
                        st.image(Image.open(io.BytesIO(msg["img"])), width=260)
                    except Exception:
                        pass
                if msg.get("snippet"):
                    # Only show expander if the snippet is not just the filename
                    if msg.get("snippet") != filename and msg.get("snippet") != msg.get("content"):
                        with st.expander("📄 File context"):
                            st.text(msg["snippet"])
                if msg.get("search_results"):
                    with st.expander("🔍 Web Search Sources", expanded=False):
                        for r in msg["search_results"]:
                            st.markdown(f"- **[{r['title']}]({r['link']})**\n  *{r['snippet']}*")
                st.markdown(msg["content"])

            if msg["role"] == "assistant":
                import hashlib
                if "tts_audio_cache" not in st.session_state:
                    st.session_state.tts_audio_cache = {}
                
                msg_content = msg["content"]
                msg_hash = hashlib.md5(msg_content.encode("utf-8")).hexdigest()
                
                export_status_placeholder = st.empty()
                col_listen, col_spacer, col_copy = st.columns([1.8, 7.6, 0.6], vertical_alignment="center")
                with col_listen:
                    if st.button("🔊 Listen", key=f"btn_listen_{idx}_{msg_hash}"):
                        st.session_state[f"play_{idx}_{msg_hash}"] = True
                        if msg_hash not in st.session_state.tts_audio_cache:
                            with st.spinner("Synthesizing..."):
                                try:
                                    voice_map = {
                                        "alice": "en-GB-SoniaNeural",
                                        "sarah": "en-US-AriaNeural",
                                        "charlie": "en-US-ChristopherNeural",
                                        "george": "en-GB-RyanNeural",
                                        "callum": "en-AU-WilliamNeural",
                                        "river": "en-US-MichelleNeural",
                                        "liam": "en-US-GuyNeural",
                                        "matilda": "en-US-JennyNeural",
                                        "will": "en-US-EricNeural",
                                        "jessica": "en-US-JennyNeural",
                                        "eric": "en-US-EricNeural",
                                        "bella": "en-US-AriaNeural",
                                        "chris": "en-US-ChristopherNeural",
                                        "brian": "en-GB-RyanNeural",
                                        "daniel": "en-US-GuyNeural",
                                        "lily": "en-US-JennyNeural",
                                        "adam": "en-US-ChristopherNeural",
                                        "bill": "en-US-EricNeural"
                                    }
                                    selected_voice = st.session_state.get("voice_response_tts_voice", "Alice")
                                    voice_id = voice_map.get(selected_voice.lower(), "en-US-AriaNeural")
                                    
                                    import asyncio
                                    import edge_tts
                                    
                                    async def _synthesize(text, voice):
                                        communicate = edge_tts.Communicate(text, voice)
                                        audio_data = b""
                                        async for chunk in communicate.stream():
                                            if chunk["type"] == "audio":
                                                audio_data += chunk["data"]
                                        return audio_data
                                    
                                    audio_bytes = asyncio.run(_synthesize(msg_content[:4000], voice_id))
                                    if audio_bytes:
                                        st.session_state.tts_audio_cache[msg_hash] = audio_bytes
                                    else:
                                        st.error("Edge-TTS synthesis failed: No audio data returned.")
                                except Exception as tts_err:
                                    st.error(f"TTS error: {tts_err}")

                with col_copy:
                    # --- Per-message Copy Button ---
                    render_copy_button(msg_content, f"copy_{idx}_{msg_hash}")
                
                if st.session_state.get(f"play_{idx}_{msg_hash}", False) and msg_hash in st.session_state.tts_audio_cache:
                    st.audio(st.session_state.tts_audio_cache[msg_hash], format="audio/mp3", autoplay=True)


# --- chat input ---
render_model_picker()
provider = st.session_state.llm_provider
if provider == "Gemini":
    model_name = st.session_state.cb_model
elif provider == "Groq":
    model_name = st.session_state.groq_model
elif provider == "BazaarLink":
    model_name = st.session_state.get("bazaarlink_model")
else:
    model_name = ""
typed = st.chat_input("Ask me anything...", accept_file="multiple", file_type=["pdf","txt","png","jpg","jpeg","webp","csv","xlsx","xls","doc","docx","pptx"])

st.markdown('<div class="voice-search-label">Voice Search</div>', unsafe_allow_html=True)
# Voice recorder — styled as a circular mic button floating next to chat input
audio_bytes = audio_recorder(text="", recording_color="#ef4444", neutral_color="#da7756", icon_name="microphone", icon_size="1x", key="voice_mic")

if audio_bytes and st.session_state.last_audio_bytes != audio_bytes:
    st.session_state.last_audio_bytes = audio_bytes
    try:
        voice_model_used, voice_model_error = get_voice_transcription_model()
        if not voice_model_used:
            st.error(f"Voice transcription error: {voice_model_error}")
            st.stop()

        with st.spinner("Transcribing..."):
            audio_model = genai.GenerativeModel(voice_model_used)
            audio_part = {"mime_type": "audio/wav", "data": audio_bytes}
            resp = audio_model.generate_content([
                "Please accurately transcribe this audio into text. Output only the transcription, nothing else.",
                audio_part
            ])
            transcribed_text = resp.text

        st.session_state.voice_prompt = transcribed_text
        if db_enabled():
            conversation_id = ensure_active_conversation()
            db_action(
                save_voice_transcription,
                st.session_state.db_user_id,
                int(conversation_id) if str(conversation_id).isdigit() else None,
                transcribed_text,
                voice_model_used,
                "audio/wav",
                audio_bytes,
            )
    except Exception as e:
        st.error(f"Voice transcription error: {e}")

if typed:
    if hasattr(typed, "files") and typed.files:
        # clear previous file types
        st.session_state.cb_file_text = ""
        st.session_state.cb_img_bytes = None
        st.session_state.cb_img_mime = None
        st.session_state.cb_df = None
        st.session_state.cb_df_name = ""
        
        for f in typed.files:
            name = f.name.lower()
            processed = False
            
            f_size = getattr(f, "size", 0)
            cache_key = f"{f.name}_{f_size}"
            
            # Check if file has already been processed in the current session
            if cache_key in st.session_state.processed_files:
                cached = st.session_state.processed_files[cache_key]
                if name.endswith('.pdf') or name.endswith('.txt') or name.endswith('.doc') or name.endswith('.docx') or name.endswith('.pptx'):
                    st.session_state.cb_file_text = cached
                elif name.endswith(('.csv', '.xlsx', '.xlsm', '.xls')):
                    st.session_state.cb_df = cached
                    st.session_state.cb_df_name = f.name
                elif name.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                    st.session_state.cb_img_bytes = cached["bytes"]
                    st.session_state.cb_img_mime = cached["mime"]
                processed = True
            else:
                if name.endswith('.pdf'):
                    try:
                        from pdfminer3.layout import LAParams
                        from pdfminer3.pdfpage import PDFPage
                        from pdfminer3.pdfinterp import PDFResourceManager, PDFPageInterpreter
                        from pdfminer3.converter import TextConverter

                        rm = PDFResourceManager()
                        sio = io.StringIO()
                        cv = TextConverter(rm, sio, laparams=LAParams())
                        pi = PDFPageInterpreter(rm, cv)
                        for pg in PDFPage.get_pages(io.BytesIO(f.read()), caching=True, check_extractable=True):
                            pi.process_page(pg)
                        st.session_state.cb_file_text = sio.getvalue()
                        cv.close(); sio.close()
                        processed = True
                        st.session_state.processed_files[cache_key] = st.session_state.cb_file_text
                    except Exception as ex:
                        st.error(f"PDF error: {ex}")
                elif name.endswith('.txt'):
                    st.session_state.cb_file_text = f.read().decode("utf-8", errors="ignore")
                    processed = True
                    st.session_state.processed_files[cache_key] = st.session_state.cb_file_text
                elif name.endswith('.doc') or name.endswith('.docx'):
                    try:
                        import docx2txt
                        st.session_state.cb_file_text = docx2txt.process(f)
                        processed = True
                        st.session_state.processed_files[cache_key] = st.session_state.cb_file_text
                    except Exception as ex:
                        st.error(f"Word document error: {ex}")
                elif name.endswith('.pptx'):
                    try:
                        st.session_state.cb_file_text = extract_pptx_text(f)
                        processed = True
                        st.session_state.processed_files[cache_key] = st.session_state.cb_file_text
                    except Exception as ex:
                        st.error(f"PowerPoint error: {ex}")
                elif name.endswith(('.csv', '.xlsx', '.xlsm', '.xls')):
                    try:
                        st.session_state.cb_df = load_dataframe_from_file(f, f.name)
                        st.session_state.cb_df_name = f.name
                        processed = True
                        st.session_state.processed_files[cache_key] = st.session_state.cb_df
                        # Auto-generate insights narrative
                        with st.spinner("AI Business Advisor is reviewing your financials..."):
                            auto_generate_insights()
                    except Exception as ex:
                        st.error(f"Dataset upload error: {ex}")
                elif name.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                    raw = f.read()
                    st.session_state.cb_img_bytes = raw
                    st.session_state.cb_img_mime = f.type
                    processed = True
                    st.session_state.processed_files[cache_key] = {"bytes": raw, "mime": f.type}
            
            if processed:
                save_attachment_record(f.name, getattr(f, "type", ""), source="chat_upload")
                # Ensure a conversation exists before recording the upload message
                ensure_active_conversation()
                
                # Determine file type for custom message and alternate turns
                file_kind = get_file_type(f.name)
                user_msg = f"Uploaded file: {f.name}"
                if file_kind == "document":
                    assistant_msg = f"📄 **{f.name}** has been successfully attached to this conversation. You can now ask questions or request an analysis of this document!"
                elif file_kind == "dataset":
                    assistant_msg = f"📊 **{f.name}** has been successfully loaded into the Data Analysis Workspace. Initial insights have been generated. What would you like to analyze?"
                elif file_kind == "image":
                    assistant_msg = f"🖼️ **{f.name}** has been successfully attached. Ask me any questions about this image!"
                else:
                    assistant_msg = f"📎 **{f.name}** has been successfully attached to this conversation."
                
                st.session_state.cb_messages.append({"role": "user", "content": user_msg, "snippet": f.name})
                st.session_state.cb_messages.append({"role": "assistant", "content": assistant_msg})
                if db_enabled():
                    conversation_id = st.session_state.active_conversation_id
                    db_action(save_message, int(conversation_id), "user", user_msg, f.name, None, None)
                    db_action(save_message, int(conversation_id), "assistant", assistant_msg, None, None, None)
                sync_active_conversation()
        
        # If user just uploaded files without typing, refresh to show previews immediately
        if not hasattr(typed, "text") or not typed.text:
            st.rerun()

text_prompt = ""
if st.session_state.voice_prompt:
    text_prompt = st.session_state.voice_prompt
    st.session_state.voice_prompt = ""
elif typed:
    if hasattr(typed, "text") and typed.text:
        text_prompt = typed.text
    elif isinstance(typed, str):
        text_prompt = typed

user_input = text_prompt or chosen

if user_input:
    # Run web search if enabled
    search_context = ""
    search_results = []
    if st.session_state.get("web_search_enabled"):
        with st.spinner(f"🔍 Searching the web for '{user_input}'..."):
            from web_search import perform_ddg_search, format_search_results_context
            search_results = perform_ddg_search(user_input, max_results=5)
            if search_results:
                search_context = format_search_results_context(search_results, user_input)
                st.session_state["last_search_results"] = search_results

    f_text  = st.session_state.cb_file_text
    i_bytes = st.session_state.cb_img_bytes
    df_ctx  = st.session_state.cb_df
    df_name = st.session_state.cb_df_name

    # Build data context string for AI
    data_ctx_str = ""
    if df_ctx is not None:
        buf = io.StringIO()
        df_ctx.info(buf=buf)
        info_str = buf.getvalue()
        desc_str = df_ctx.describe(include='all').to_string()
        sample_str = df_ctx.head(20).to_string(index=False)
        data_ctx_str = (
            f"\n\n[Uploaded Dataset: {df_name}]\n"
            f"Shape: {df_ctx.shape[0]} rows x {df_ctx.shape[1]} columns\n"
            f"Columns: {', '.join(df_ctx.columns.tolist())}\n\n"
            f"Data Info:\n{info_str}\n"
            f"Statistical Summary:\n{desc_str}\n\n"
            f"First 20 rows:\n{sample_str}"
        )

    snippet = (f_text[:400] if f_text else "") or (data_ctx_str[:400] if data_ctx_str else None)
    st.session_state.cb_messages.append({
        "role": "user", "content": user_input,
        "img": i_bytes,
        "snippet": snippet,
        "search_results": search_results if search_results else None
    })
    conversation_id = ensure_active_conversation()
    if db_enabled():
        db_action(
            save_message,
            int(conversation_id),
            "user",
            user_input,
            snippet,
            i_bytes,
            st.session_state.cb_img_mime,
        )
    sync_active_conversation()

    with chat_box:
        with st.chat_message("user"):
            if i_bytes:
                try:
                    from PIL import Image
                    st.image(Image.open(io.BytesIO(i_bytes)), width=260)
                except Exception:
                    pass
            if f_text:
                with st.expander("📄 File context"):
                    st.text(f_text[:400])
            if df_ctx is not None:
                with st.expander(f"📊 Data context: {df_name}"):
                    st.dataframe(df_ctx.head(5), use_container_width=True)
            st.markdown(user_input)

        with st.chat_message("assistant"):
            # Reuse the cached Google credentials from session state
            google_creds = st.session_state.get("google_credentials")

            access_token = None
            if google_creds:
                from backend.google_service import get_valid_token
                access_token, updated_creds = get_valid_token(st.session_state.db_user_id, google_creds)
                if updated_creds:
                    st.session_state.google_credentials = updated_creds

            sys_prompt = (
                "You are a powerful, friendly, and general-purpose AI assistant. "
                "You can perform all tasks that an AI chatbot can do, answer any question, "
                "analyse tabular data (CSV/Excel), identify trends, summarise datasets, "
                "interpret statistics, analyse documents and images, and explain complex topics. "
                "When given dataset context, reference specific columns, values, and statistics in your answers. "
                "Format responses clearly with headings, bullet points, and tables where useful. "
                "CRITICAL GMAIL/CALENDAR RULES: When using Gmail or Calendar tools, you MUST pass the COMPLETE, "
                "FULL email body text into the tool 'body' argument — never truncate, summarize, or shorten it. "
                "After the tool executes successfully, your response to the user must ONLY be a brief confirmation "
                "like: 'Done! I\'ve saved the email draft to your Gmail Drafts folder.' — do NOT reproduce, "
                "rewrite, or paraphrase the email body in the chat. The email content is already saved in Gmail exactly as written."
            )
            # Query vector database for document context if active conversation exists
            vector_context = ""
            if db_enabled() and st.session_state.get("active_conversation_id"):
                try:
                    from backend.vector_service import query_relevant_context
                    vector_context = query_relevant_context(
                        st.session_state.active_conversation_id,
                        user_input
                    )
                except Exception as ve:
                    print(f"Error querying vector DB in Streamlit: {ve}")

            if vector_context:
                file_ctx = f"\n\n[Relevant Context from uploaded files]:\n{vector_context}"
            else:
                file_ctx = f"\n\n[Document content (first 3000 chars)]:\n{f_text[:3000]}" if f_text else ""
            
            current_text = f"{file_ctx}{data_ctx_str}\n\nUser: {user_input}" if (file_ctx or data_ctx_str) else user_input

            if 'search_context' in locals() and search_context:
                sys_prompt += (
                    "\n\nYou have access to real-time search results. Use the web search results "
                    "below to answer the user's request. You must cite your sources using inline links (e.g., [Title](URL)) "
                    "or list them at the end of your response under a 'Sources' section."
                )
                current_text = f"{search_context}\n\n{current_text}"

            is_vision = is_vision_model(provider, model_name)
            if i_bytes and not is_vision:
                st.warning(f"⚠️ The selected model **{model_name}** ({provider}) does not support image analysis. The query will be processed as text-only. Switch to a vision-supporting model to analyze the image.")

            try:
                # Setup Google tools schema for OpenAI-compatible paths
                tools_list = []
                # Disable tools for BazaarLink's auto:free model as it does not reliably support function calling
                if access_token and not (provider == "BazaarLink" and model_name == "auto:free"):
                    tools_list = [
                        {
                            "type": "function",
                            "function": {
                                "name": "draft_gmail_email",
                                "description": "Save a full email as a draft in the user's Gmail Drafts folder. Use this when the user says 'draft', 'save as draft', or doesn't explicitly ask to send. Pass the COMPLETE, UNTRUNCATED email body.",
                                "parameters": {
                                    "type": "object",
                                    "properties": {
                                        "recipient": {"type": "string", "description": "The recipient email address."},
                                        "subject": {"type": "string", "description": "The email subject line."},
                                        "body": {"type": "string", "description": "The COMPLETE, FULL plain-text body of the email. Do not truncate or summarize — include every sentence, bullet point, paragraph, and signature exactly as written."}
                                    },
                                    "required": ["recipient", "subject", "body"]
                                }
                            }
                        },
                        {
                            "type": "function",
                            "function": {
                                "name": "send_gmail_email",
                                "description": "Send an email directly and immediately via Gmail (not saved as draft). Use this when the user explicitly says 'send' the email. Pass the COMPLETE, UNTRUNCATED email body.",
                                "parameters": {
                                    "type": "object",
                                    "properties": {
                                        "recipient": {"type": "string", "description": "The recipient email address."},
                                        "subject": {"type": "string", "description": "The email subject line."},
                                        "body": {"type": "string", "description": "The COMPLETE, FULL plain-text body of the email. Do not truncate or summarize — include every sentence, bullet point, paragraph, and signature exactly as written."}
                                    },
                                    "required": ["recipient", "subject", "body"]
                                }
                            }
                        },
                        {
                            "type": "function",
                            "function": {
                                "name": "create_calendar_event",
                                "description": "Create a new event in Google Calendar.",
                                "parameters": {
                                    "type": "object",
                                    "properties": {
                                        "summary": {"type": "string", "description": "The event title."},
                                        "start_time": {"type": "string", "description": "Start time in ISO 8601 format (e.g. '2026-06-15T09:00:00')."},
                                        "end_time": {"type": "string", "description": "End time in ISO 8601 format (e.g. '2026-06-15T10:00:00')."},
                                        "description": {"type": "string", "description": "Description of the event (optional)."},
                                        "location": {"type": "string", "description": "Location of the event (optional)."}
                                    },
                                    "required": ["summary", "start_time", "end_time"]
                                }
                            }
                        }
                    ]

                if provider in ("Groq", "BazaarLink"):
                    OpenAI = import_openai_client()
                    if OpenAI is None:
                        st.error("Install the OpenAI package first: pip install openai")
                        st.stop()

                    if provider == "Groq":
                        if not GROQ_API_KEY or GROQ_API_KEY == "YOUR_GROQ_API_KEY_HERE":
                            st.error("Add your Groq API key in App.py to use Groq.")
                            st.stop()
                        client = OpenAI(
                            base_url="https://api.groq.com/openai/v1",
                            api_key=GROQ_API_KEY,
                            timeout=30.0
                        )
                    else: # BazaarLink
                        if not BAZAARLINK_API_KEY:
                            st.error("Add your BazaarLink API key in secrets.toml to use BazaarLink.")
                            st.stop()
                        client = OpenAI(
                            base_url=BAZAARLINK_BASE_URL,
                            api_key=BAZAARLINK_API_KEY,
                            default_headers=BAZAARLINK_HEADERS,
                            timeout=30.0
                        )

                    if not model_name:
                        st.error(f"No compatible {provider} model was found for this API key.")
                        st.stop()

                    messages = build_openai_messages(
                        sys_prompt,
                        st.session_state.cb_messages[:-1],
                        current_text,
                        i_bytes,
                        st.session_state.cb_img_mime,
                        is_vision=is_vision
                    )

                    ph = st.empty()
                    answer = ""

                    if tools_list:
                        response = client.chat.completions.create(
                            model=model_name,
                            messages=messages,
                            tools=tools_list,
                            tool_choice="auto"
                        )
                        response_message = response.choices[0].message
                        if response_message.tool_calls:
                            import json
                            from backend.google_service import draft_email, send_email, create_event

                            tool_confirmations = []
                            messages.append(response_message)
                            for tool_call in response_message.tool_calls:
                                args = json.loads(tool_call.function.arguments)
                                fn_name = tool_call.function.name
                                if fn_name == "draft_gmail_email":
                                    body_text = args.get("body", "")
                                    res = draft_email(args.get("recipient", ""), args.get("subject", ""), body_text, access_token)
                                    if res.startswith("Successfully"):
                                        tool_confirmations.append(
                                            f"✅ **Email draft saved to Gmail Drafts!**\n\n"
                                            f"**To:** {args.get('recipient', '')}\n"
                                            f"**Subject:** {args.get('subject', '')}\n\n"
                                            f"---\n\n{body_text}"
                                        )
                                    else:
                                        tool_confirmations.append(f"❌ {res}")
                                elif fn_name == "send_gmail_email":
                                    body_text = args.get("body", "")
                                    res = send_email(args.get("recipient", ""), args.get("subject", ""), body_text, access_token)
                                    if res.startswith("Email successfully"):
                                        tool_confirmations.append(
                                            f"✅ **Email sent successfully!**\n\n"
                                            f"**To:** {args.get('recipient', '')}\n"
                                            f"**Subject:** {args.get('subject', '')}\n\n"
                                            f"---\n\n{body_text}"
                                        )
                                    else:
                                        tool_confirmations.append(f"❌ {res}")
                                elif fn_name == "create_calendar_event":
                                    res = create_event(
                                        args.get("summary", ""),
                                        args.get("start_time", ""),
                                        args.get("end_time", ""),
                                        args.get("description", ""),
                                        args.get("location", ""),
                                        access_token
                                    )
                                    tool_confirmations.append(f"✅ {res}")
                                else:
                                    res = "Unknown function call"
                                    tool_confirmations.append(res)

                                messages.append({
                                    "role": "tool",
                                    "tool_call_id": tool_call.id,
                                    "name": fn_name,
                                    "content": res
                                })

                            # Show the tool results directly — do NOT ask LLM to rewrite the email
                            answer = "\n\n".join(tool_confirmations)
                            ph.markdown(answer)
                        else:
                            answer = response_message.content or ""
                            ph.markdown(answer)
                    else:
                        stream = client.chat.completions.create(
                            model=model_name,
                            messages=messages,
                            stream=True,
                        )
                        for chunk in stream:
                            delta = chunk.choices[0].delta.content
                            if delta:
                                answer += delta
                                ph.markdown(answer + "◼")
                        ph.markdown(answer)

                    st.session_state.cb_messages.append({
                        "role": "assistant", 
                        "content": answer,
                        "search_results": search_results if ('search_results' in locals() and search_results) else None
                    })
                    if db_enabled():
                        db_action(save_message, int(conversation_id), "assistant", answer)
                        db_action(save_api_log, int(conversation_id), provider, model_name, current_text, answer)
                    sync_active_conversation()
                    # Trigger UI refresh for recent chats after assistant response
                    st.session_state.just_responded = True
                    st.rerun()
                    st.stop()

                if not GEMINI_API_KEY or GEMINI_API_KEY == "YOUR_GEMINI_API_KEY_HERE":
                    st.error("Add your Gemini API key in App.py to use Gemini.")
                    st.stop()
                if not model_name:
                    st.error("No compatible Gemini model was found for this API key.")
                    st.stop()

                genai = import_genai_module()
                if genai is None:
                    st.error("Install google-generativeai first: pip install google-generativeai")
                    st.stop()

                # Prepare Gemini tools
                def open_local_browser_tab(url: str) -> str:
                    """Opens a web page in a new browser tab.
                    Use this tool when the user asks to open a website, browse a URL locally, or open a link in a tab.

                    Args:
                        url: The exact HTTP/HTTPS URL of the website to open.
                    """
                    import webbrowser
                    try:
                        clean_url = url.strip().strip("'\"")
                        webbrowser.open(clean_url)
                        return f"Successfully opened new tab for: {clean_url}"
                    except Exception as e:
                        return f"Failed to open browser tab: {e}"

                def browse_webpage(url: str) -> str:
                    """Fetches and reads the text content of a webpage so you can answer questions about it.
                    Use this tool when the user asks to read, analyze, search, or summarize the contents of a specific URL.

                    Args:
                        url: The HTTP/HTTPS URL of the webpage to read.
                    """
                    import urllib.request
                    import urllib.parse
                    from bs4 import BeautifulSoup
                    try:
                        clean_url = url.strip().strip("'\"")
                        req = urllib.request.Request(
                            clean_url,
                            headers={
                                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
                            }
                        )
                        with urllib.request.urlopen(req, timeout=10) as response:
                            html = response.read()
                        soup = BeautifulSoup(html, "html.parser")
                        for s in soup(["script", "style", "meta", "noscript", "header", "footer", "nav"]):
                            s.decompose()
                        text = soup.get_text(separator=" ", strip=True)
                        truncated_text = text[:4000]
                        if len(text) > 4000:
                            truncated_text += "\n\n[Content truncated for length limit]"
                        return f"Successfully fetched content from {clean_url}:\n\n{truncated_text}"
                    except Exception as e:
                        return f"Failed to read webpage content from {url}: {e}"

                tools = [open_local_browser_tab, browse_webpage]
                _gemini_tool_results = {}  # store tool results keyed by fn name for display
                if access_token:
                    from backend.google_service import draft_email, send_email, create_event

                    def draft_gmail_email(recipient: str, subject: str, body: str) -> str:
                        """Save a full email as a draft in Gmail Drafts. Use when user says 'draft'. Pass the COMPLETE, FULL email body — never truncate.

                        Args:
                            recipient: The recipient email address.
                            subject: The email subject line.
                            body: The COMPLETE, FULL plain-text body of the email. Include every sentence and bullet point.
                        """
                        result = draft_email(recipient, subject, body, access_token)
                        _gemini_tool_results['draft_gmail_email'] = {
                            'recipient': recipient, 'subject': subject, 'body': body, 'result': result
                        }
                        return result

                    def send_gmail_email(recipient: str, subject: str, body: str) -> str:
                        """Send an email directly via Gmail. Use when user says 'send'. Pass the COMPLETE, FULL email body — never truncate.

                        Args:
                            recipient: The recipient email address.
                            subject: The email subject line.
                            body: The COMPLETE, FULL plain-text body of the email. Include every sentence and bullet point.
                        """
                        result = send_email(recipient, subject, body, access_token)
                        _gemini_tool_results['send_gmail_email'] = {
                            'recipient': recipient, 'subject': subject, 'body': body, 'result': result
                        }
                        return result

                    def create_calendar_event(summary: str, start_time: str, end_time: str, description: str = "", location: str = "") -> str:
                        """Create a new event in Google Calendar.

                        Args:
                            summary: The event title.
                            start_time: Start time in ISO 8601 format (e.g. '2026-06-15T09:00:00').
                            end_time: End time in ISO 8601 format (e.g. '2026-06-15T10:00:00').
                            description: Description of the event (optional).
                            location: Location of the event (optional).
                        """
                        return create_event(summary, start_time, end_time, description, location, access_token)

                    tools.extend([draft_gmail_email, send_gmail_email, create_calendar_event])

                model = genai.GenerativeModel(model_name, system_instruction=sys_prompt, tools=tools)
                
                history = []
                # Add previous messages to context
                for m in st.session_state.cb_messages[:-1]: # exclude current
                    role = "user" if m["role"] == "user" else "model"
                    history.append({"role": role, "parts": [m["content"]]})
                
                chat = model.start_chat(history=history, enable_automatic_function_calling=bool(tools))
                
                # Build current prompt parts
                current_parts = [current_text]
                
                # Add image content if present
                if i_bytes:
                    mime = st.session_state.cb_img_mime or "image/jpeg"
                    current_parts.append({"mime_type": mime, "data": i_bytes})

                if tools or i_bytes:
                    resp = chat.send_message(current_parts)
                    # Check if any Gmail tool was called — show the full original email body instead of LLM rewrite
                    if _gemini_tool_results:
                        confirmation_parts = []
                        for fn_key, info in _gemini_tool_results.items():
                            res = info.get('result', '')
                            body_text = info.get('body', '')
                            if fn_key == 'draft_gmail_email' and res.startswith('Successfully'):
                                confirmation_parts.append(
                                    f"✅ **Email draft saved to Gmail Drafts!**\n\n"
                                    f"**To:** {info.get('recipient', '')}\n"
                                    f"**Subject:** {info.get('subject', '')}\n\n"
                                    f"---\n\n{body_text}"
                                )
                            elif fn_key == 'send_gmail_email' and res.startswith('Email successfully'):
                                confirmation_parts.append(
                                    f"✅ **Email sent successfully!**\n\n"
                                    f"**To:** {info.get('recipient', '')}\n"
                                    f"**Subject:** {info.get('subject', '')}\n\n"
                                    f"---\n\n{body_text}"
                                )
                            else:
                                confirmation_parts.append(res)
                        answer = "\n\n".join(confirmation_parts) if confirmation_parts else resp.text
                    else:
                        answer = resp.text
                    st.markdown(answer)
                    st.session_state.cb_messages.append({
                        "role": "assistant", 
                        "content": answer,
                        "search_results": search_results if ('search_results' in locals() and search_results) else None
                    })
                    if db_enabled():
                        db_action(save_message, int(conversation_id), "assistant", answer)
                        db_action(save_api_log, int(conversation_id), provider, model_name, current_text, answer)
                    sync_active_conversation()
                    # Trigger UI refresh for recent chats after assistant response
                    st.session_state.just_responded = True
                    st.rerun()
                    st.stop()
                else:
                    stream = chat.send_message(current_parts, stream=True)
                    ph = st.empty()
                    out = ""
                    for chunk in stream:
                        if chunk.text:
                            out += chunk.text
                            ph.markdown(out + "◼")
                    ph.markdown(out)
                    st.session_state.cb_messages.append({
                        "role": "assistant", 
                        "content": out,
                        "search_results": search_results if ('search_results' in locals() and search_results) else None
                    })
                    if db_enabled():
                        db_action(save_message, int(conversation_id), "assistant", out)
                        db_action(save_api_log, int(conversation_id), provider, model_name, current_text, out)
                    sync_active_conversation()
                    # Trigger UI refresh for recent chats after assistant response
                    st.session_state.just_responded = True
                    st.rerun()
                    st.stop()
            except Exception as ex:
                if db_enabled():
                    db_action(
                        save_api_log,
                        int(conversation_id),
                        provider,
                        model_name,
                        current_text if "current_text" in locals() else user_input,
                        "",
                        "error",
                        str(ex),
                    )
                st.error(f"Generation error: {ex}")
record_profiler_checkpoint("Chat View & Input")